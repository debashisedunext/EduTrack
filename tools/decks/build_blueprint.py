#!/usr/bin/env python3
"""EduTrack — Product Blueprint deck, generated from Ticketing-System-Blueprint.md"""
import sys, pathlib
sys.path.insert(0, str(pathlib.Path(__file__).parent))
from deck_lib import *
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

new_deck()
FT = "Product Blueprint v1.2"
n = [0]
def pg(s):
    n[0] += 1
    footer(s, n[0], FT)

def hero(title, kicker, sub):
    s = slide(SURFACE)
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SW), Inches(SH))
    b.fill.gradient(); b.fill.gradient_angle = 45.0
    st = b.fill.gradient_stops
    st[0].color.rgb = RGBColor(0x4F,0x46,0xE5); st[0].position = 0.0
    st[1].color.rgb = RGBColor(0x1E,0x1B,0x4B); st[1].position = 1.0
    b.line.fill.background(); b.shadow.inherit = False
    txt(s, ML+0.3, 2.05, CW, 0.4, kicker, size=15, bold=True, color=RGBColor(0xA5,0xB4,0xFC))
    txt(s, ML+0.3, 2.55, CW, 1.6, title, size=42, bold=True, color=WHITE, line_spacing=1.08)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML+0.3), Inches(4.45), Inches(1.5), Inches(0.045))
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0xA5,0xB4,0xFC)
    ln.line.fill.background(); ln.shadow.inherit = False
    txt(s, ML+0.3, 4.78, CW-1.0, 1.0, sub, size=15.5, color=RGBColor(0xC7,0xD2,0xFE), line_spacing=1.35)
    return s

# ---------------------------------------------------------------- 1
hero("Organisation Task &\nClient Ticketing System",
     "EDUTRACK",
     "34 screens · 6 roles · ~28 tables · 20 dashboard widgets · 18 reports\n"
     "Complete product, architecture & UI blueprint")

# ---------------------------------------------------------------- 2 exec
s = slide()
y = header(s, "What EduTrack merges",
           "Two things most organisations run as separate products, unified on one work item")
rows = [
    ["Project + task assignment, sprints, velocity", ("Jira Software · Azure DevOps", False, TEXT2)],
    ["Client ticket intake, SLA, escalation", ("Zoho Desk · Freshdesk · Zendesk", False, TEXT2)],
    ["Quick update / personal work queue", ("Linear · ClickUp \"My Work\"", False, TEXT2)],
    ["Chat threaded on the work item", ("Slack + Jira — unified here", False, TEXT2)],
    ["Analytical dashboards with drill-down", ("Jira Dashboards · Freshdesk Analytics", False, TEXT2)],
]
table(s, ML, y, CW, ["Capability", "Reference product we borrow from"], rows,
      [7.0, 4.93], row_h=0.48, head_h=0.42, fsize=13)
card(s, ML, y+3.05, CW, 1.5, fill=PRIMARY_SOFT, line=None)
txt(s, ML+0.35, y+3.28, CW-0.7, 1.1,
    "An internal, multi-project task and client ticketing platform.\n\n"
    "The unification is the point: a ticket carries its project work, its client relationship, its SLA "
    "and its conversation in one place — so nobody reconciles two systems to answer one question.",
    size=13.5, color=PRIMARY, line_spacing=1.35)
pg(s)

# ---------------------------------------------------------------- 3 differentiators
s = slide()
y = header(s, "Six things that make it more than a Jira clone", "")
items = [
    ("Cycle-based reopen model", "Every reopen creates a new cycle with its own start date, assignee, planned close date and effort. Nothing is overwritten, ever."),
    ("Immutable history", "Resources update state, never history. Append-only at the database level, not just in the UI."),
    ("Automatic criticality escalation", "A ticket crossing its Planned Close Date is auto-promoted to Critical and pushed to the reporting manager."),
    ("Chat attached to the ticket", "The manager's \"what's the status?\" becomes a first-class, trackable, timed event — not a side channel."),
    ("The Workflow Ribbon", "A live journey strip showing the full Support → PM → Dev → QA → Deploy → Close chain, where the ticket is, how many times it bounced, and each person's effort."),
    ("Client-aware from the first field", "Client master with Excel import, client dropdown on every ticket, client-visible vs internal comments, mail threading back as comments."),
]
cy = y
for i, (t, b) in enumerate(items):
    x = ML if i % 2 == 0 else ML + 6.15
    yy = cy + (i // 2) * 1.42
    card(s, x, yy, 5.83, 1.28)
    num = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x+0.28), Inches(yy+0.26), Inches(0.38), Inches(0.38))
    num.fill.solid(); num.fill.fore_color.rgb = PRIMARY
    num.line.fill.background(); num.shadow.inherit = False
    txt(s, x+0.28, yy+0.33, 0.38, 0.28, str(i+1), size=12.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, x+0.82, yy+0.24, 4.8, 0.32, t, size=14.5, bold=True)
    txt(s, x+0.82, yy+0.62, 4.8, 0.6, b, size=11, color=TEXT2, line_spacing=1.22)
pg(s)

# ---------------------------------------------------------------- 4 roles
s = slide()
y = header(s, "Six roles, extensible via a role master",
           "QA and Deployment exist because the ribbon needs a receiving role that owns each stage")
rows = [
    [("Create/edit resources, roles, reporting manager", True, TEXT), ("✔", True, SUCCESS), "—", "—", "—", "—", "—"],
    [("Create/edit projects, map resources", True, TEXT), ("✔", True, SUCCESS), ("✔ own", True, SUCCESS), "—", "—", "—", "—"],
    [("Create ticket", True, TEXT), ("✔", True, SUCCESS), ("✔", True, SUCCESS), ("✔", True, SUCCESS), ("✔", True, SUCCESS), ("✔", True, SUCCESS), ("✔", True, SUCCESS)],
    [("Assign / reassign ticket", True, TEXT), ("✔", True, SUCCESS), ("✔", True, SUCCESS), ("✔ own", True, SUCCESS), "—", "—", "—"],
    [("Hand off to next stage · send back", True, TEXT), ("✔", True, SUCCESS), ("✔", True, SUCCESS), ("✔", True, SUCCESS), ("✔", True, SUCCESS), ("✔", True, SUCCESS), ("✔", True, SUCCESS)],
    [("Skip a stage · force-move ribbon", True, TEXT), ("✔", True, SUCCESS), ("✔", True, SUCCESS), "—", "—", "—", "—"],
    [("View all tickets", True, TEXT), ("✔", True, SUCCESS), ("proj", False, TEXT2), ("proj", False, TEXT2), ("mine", False, TEXT2), ("mine", False, TEXT2), ("mine", False, TEXT2)],
    [("Close / reopen ticket", True, TEXT), ("✔", True, SUCCESS), ("✔", True, SUCCESS), ("✔", True, SUCCESS), "—", "—", "—"],
    [("Edit or delete history / ribbon", True, TEXT), ("✘", True, DANGER), ("✘", True, DANGER), ("✘", True, DANGER), ("✘", True, DANGER), ("✘", True, DANGER), ("✘", True, DANGER)],
]
table(s, ML, y, CW, ["Capability", "Admin", "PM", "Support", "Dev", "QA", "Deploy"],
      rows, [5.0, 1.16, 1.16, 1.16, 1.16, 1.16, 1.13], row_h=0.35, head_h=0.38, fsize=11)
card(s, ML, y+3.72, CW, 0.72, fill=RGBColor(0xFE,0xF2,0xF2), line=None)
txt(s, ML+0.32, y+3.88, CW-0.64, 0.5,
    "Nobody — including Admin — can edit or delete history or the ribbon. That row is ✘ across all six columns by design.",
    size=13, color=RGBColor(0x99,0x1B,0x1B), line_spacing=1.3)
pg(s)

# ---------------------------------------------------------------- 5 scopes
s = slide()
y = header(s, "Two orthogonal scopes decide visibility on every query", "")
card(s, ML, y, 5.8, 1.55, fill=PRIMARY_SOFT, line=None)
txt(s, ML+0.32, y+0.24, 5.2, 0.3, "ROLE SCOPE", size=11, bold=True, color=PRIMARY)
txt(s, ML+0.32, y+0.62, 5.2, 0.8, "What the role is allowed to do.\nCreate, assign, close, configure masters.",
    size=13.5, color=PRIMARY, line_spacing=1.3)
card(s, ML+6.15, y, 5.78, 1.55, fill=RGBColor(0xEC,0xFD,0xF5), line=None)
txt(s, ML+6.47, y+0.24, 5.2, 0.3, "ROW SCOPE", size=11, bold=True, color=RGBColor(0x06,0x5F,0x46))
txt(s, ML+6.47, y+0.62, 5.2, 0.8,
    "assignee = me  OR  reported_by = me\nOR project_id IN my_projects  OR  assignee IN my_reportees",
    size=12, color=RGBColor(0x06,0x5F,0x46), line_spacing=1.3, font=MONO)
card(s, ML, y+1.85, CW, 1.35, fill=RGBColor(0x1E,0x1B,0x33), line=None)
txt(s, ML+0.35, y+2.05, CW-0.7, 0.3, "THE SINGLE MOST IMPORTANT SECURITY RULE", size=11, bold=True, color=RGBColor(0xA5,0xB4,0xFC))
txt(s, ML+0.35, y+2.45, CW-0.7, 0.7,
    "A Developer's ticket list is always forcibly filtered by assigned_to = current_user_id, injected by a "
    "server-side scope guard — never by a front-end filter.",
    size=14, color=WHITE, line_spacing=1.3)
card(s, ML, y+3.5, CW, 1.0, fill=SURFACE)
txt(s, ML+0.35, y+3.68, CW-0.7, 0.3, "GOLDEN RULE OF THE RIBBON", size=11, bold=True, color=TEXT2)
txt(s, ML+0.35, y+4.02, CW-0.7, 0.4,
    "Only the current stage owner — plus PM and Admin — can move a ticket to the next stage. "
    "A Developer cannot push a ticket into Deployment while it is sitting with QA.",
    size=13.5, color=TEXT, line_spacing=1.3)
pg(s)

# ---------------------------------------------------------------- 6 stage vs status
s = slide()
y = header(s, "Two layers: stage and status",
           "Orthogonal, and both are needed — the ribbon alone cannot tell you whether work is moving")
rows = [
    [("Answers", True, TEXT), "Which team owns it right now?", "Is work moving, or blocked?"],
    [("Values", True, TEXT), ("Intake · Triage · Development · QA · Deployment · Verification · Sign-off · Closed", False, TEXT),
     ("New · In Progress · On Hold · Awaiting Info · Rework · Resolved · Closed · Reopened", False, TEXT)],
    [("Changes on", True, TEXT), "A handoff between teams", "A day-to-day update by the owner"],
    [("Owner", True, TEXT), "Stage owner role", "Assignee"],
]
table(s, ML, y, CW, ["", "Stage  (the ribbon)", "Status"], rows,
      [1.6, 5.2, 5.13], row_h=0.72, head_h=0.42, fsize=12.5)
card(s, ML, y+3.42, CW, 1.25, fill=PRIMARY_SOFT, line=None)
txt(s, ML+0.35, y+3.62, CW-0.7, 0.9,
    "A ticket can be In Progress in the QA stage, or On Hold in the Deployment stage.\n\n"
    "Without both, you cannot distinguish \"sitting in QA because Anil is busy\" from "
    "\"sitting in QA because the test environment is down\".",
    size=13.5, color=PRIMARY, line_spacing=1.35)
pg(s)

# ---------------------------------------------------------------- 7 flow
s = slide()
y = header(s, "End-to-end operational flow", "Seven phases, from intake to reopen")
steps = [
    ("1", "INTAKE", "Support Desk / PM / client raises a ticket, or the email parser does. ID auto-generated: CRM-26-00347. SLA policy resolves the planned close date.", A_C),
    ("2", "TRIAGE", "PM or Support validates, sets level, assigns. Cycle #1 stamped. Real-time popup + bell + email to the assignee within ~1 second.", B_C),
    ("3", "EXECUTION", "Moves along the ribbon. At every hop: effort attributed to that resource, time-in-stage stamped, an append-only transition row written.", C_C),
    ("4", "MONITORING", "Scheduler every 15 minutes: past planned close → Critical + alert to reporting manager. 80% elapsed → pre-breach warning. No update in N days → nudge.", D_C),
    ("5", "STATUS CHECK", "Manager clicks Ask Status → templated message into the ticket's chat thread. Response time is recorded as a metric.", A_C),
    ("6", "CLOSURE", "Resource marks Resolved → PM/Support verifies → Closed. Actual close date stamped, cycle effort frozen, optional CSAT.", B_C),
    ("7", "REOPEN", "Reason mandatory. Cycle #1 sealed read-only forever. Cycle #2 opens with fresh dates and its own ribbon. Total effort = Σ all cycles.", C_C),
]
cy = y
for num, t, b, col in steps:
    card(s, ML, cy, CW, 0.58)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML), Inches(cy+0.08), Inches(0.055), Inches(0.42))
    bar.fill.solid(); bar.fill.fore_color.rgb = col
    bar.line.fill.background(); bar.shadow.inherit = False
    txt(s, ML+0.28, cy+0.16, 0.3, 0.3, num, size=13, bold=True, color=col)
    txt(s, ML+0.62, cy+0.16, 1.75, 0.3, t, size=12.5, bold=True, color=TEXT)
    txt(s, ML+2.45, cy+0.14, CW-2.75, 0.42, b, size=11.5, color=TEXT2, line_spacing=1.15)
    cy += 0.66
pg(s)

# ---------------------------------------------------------------- 8 history model
s = slide()
y = header(s, "The reopen & history model",
           "Three layers, each with a different job — this is where most in-house ticketing tools fail")
rows = [
    [("Current state", True, TEXT), ("tickets", False, TEXT), "Fast reads, list views, dashboards", ("Updatable", True, INFO)],
    [("Cycle snapshot", True, TEXT), ("ticket_cycles", False, TEXT), "One row per open → close round", ("Insert; sealed on close", True, WARNING)],
    [("Change trail", True, TEXT), ("ticket_history", False, TEXT), "Every single field change, who and when", ("APPEND-ONLY", True, DANGER)],
    [("Effort trail", True, TEXT), ("ticket_effort_logs", False, TEXT), "Every hour logged, by whom, on which date", ("APPEND-ONLY", True, DANGER)],
    [("Stage trail", True, TEXT), ("ticket_stage_transitions", False, TEXT), "Every hop — this table IS the ribbon", ("APPEND-ONLY", True, DANGER)],
]
table(s, ML, y, CW, ["Layer", "Table", "Purpose", "Mutability"], rows,
      [2.0, 3.0, 4.6, 2.33], row_h=0.48, head_h=0.42, fsize=12)
txt(s, ML, y+3.1, CW, 0.3, "WHAT A REOPEN ACTUALLY DOES", size=11, bold=True, color=TEXT2)
code_block(s, ML, y+3.46, CW, 1.2,
           "BEGIN;\n"
           "  UPDATE ticket_cycles SET is_sealed = true WHERE ticket_id = :id AND cycle_no = :current;\n"
           "  INSERT INTO ticket_cycles (…, cycle_no, start_date, assigned_to, planned_close_date, reopen_reason)\n"
           "  UPDATE tickets SET status='REOPENED', current_cycle_no = :current+1, reopen_count = reopen_count+1;\n"
           "COMMIT;        -- effort from cycle 1 is never touched", size=10.5)
pg(s)

# ---------------------------------------------------------------- 9 four guards
s = slide()
y = header(s, "How history is made tamper-proof", "Four independent guards — defence in depth")
guards = [
    ("Application", "No service method exists that issues UPDATE or DELETE on the protected tables. Only insert() is exposed.", A_C),
    ("Database", "The app's DB role holds INSERT and SELECT grants only. A BEFORE UPDATE OR DELETE trigger raises an exception.", B_C),
    ("API", "No PUT, PATCH or DELETE routes are registered. A correction is a new compensating entry — exactly like an accounting reversal.", C_C),
    ("Integrity", "Each row stores prev_hash + row_hash as a SHA-256 chain. A nightly job verifies it; a break raises an admin alert.", D_C),
]
cwid = (CW - 0.75) / 4
for i, (t, b, col) in enumerate(guards):
    x = ML + i * (cwid + 0.25)
    card(s, x, y, cwid, 2.9)
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(cwid), Inches(0.07))
    top.fill.solid(); top.fill.fore_color.rgb = col
    top.line.fill.background(); top.shadow.inherit = False
    txt(s, x+0.26, y+0.34, cwid-0.52, 0.3, str(i+1).zfill(2), size=12, bold=True, color=col)
    txt(s, x+0.26, y+0.72, cwid-0.52, 0.35, t, size=16, bold=True)
    txt(s, x+0.26, y+1.2, cwid-0.52, 1.5, b, size=11.5, color=TEXT2, line_spacing=1.28)
card(s, ML, y+3.15, CW, 1.35, fill=RGBColor(0x1E,0x1B,0x33), line=None)
txt(s, ML+0.35, y+3.35, CW-0.7, 0.3, "WHY FOUR AND NOT ONE", size=11, bold=True, color=RGBColor(0xA5,0xB4,0xFC))
txt(s, ML+0.35, y+3.72, CW-0.7, 0.7,
    "The hash chain makes even direct database tampering detectable. For a Developer, the history grid renders "
    "with no edit or delete affordance at all — and the API would reject it even if the DOM were manipulated.",
    size=13.5, color=WHITE, line_spacing=1.3)
pg(s)

# ---------------------------------------------------------------- 10 ribbon stages
s = slide()
y = header(s, "The Workflow Ribbon — eight stages",
           "A live journey strip pinned to the top of every ticket")
rows = [
    [("1", True, TEXT2), ("Intake", True, TEXT), "Support Desk", "Ticket created", "Details validated, project tagged", ("2 h", False, TEXT2)],
    [("2", True, TEXT2), ("Triage / Planning", True, TEXT), "PM", "Support hands off", "Level set + developer assigned", ("4 h", False, TEXT2)],
    [("3", True, TEXT2), ("Development", True, TEXT), "Developer", "PM assigns", "Dev marks Ready for QA", ("per SLA", False, TEXT2)],
    [("4", True, TEXT2), ("QA / Testing", True, TEXT), "QA", "Dev hands off", "QA marks Pass or Fail", ("8 h", False, TEXT2)],
    [("5", True, TEXT2), ("Deployment", True, TEXT), "Deployment / DevOps", "QA passes", "Deploy marked done", ("4 h", False, TEXT2)],
    [("6", True, TEXT2), ("Verification", True, TEXT), "Developer", "Deploy hands back", "Dev confirms in target env", ("4 h", False, TEXT2)],
    [("7", True, TEXT2), ("Sign-off", True, TEXT), "PM / Support Desk", "Dev marks complete", "PM/Support accepts", ("8 h", False, TEXT2)],
    [("8", True, TEXT2), ("Closed", True, TEXT), "—", "Sign-off accepted", "terminal", ("—", False, TEXT2)],
]
table(s, ML, y, CW, ["#", "Stage", "Owner role", "Enters when", "Leaves when", "SLA"], rows,
      [0.5, 2.3, 2.3, 2.5, 3.2, 1.13], row_h=0.37, head_h=0.4, fsize=11)
card(s, ML, y+3.62, CW, 0.85, fill=PRIMARY_SOFT, line=None)
txt(s, ML+0.35, y+3.8, CW-0.7, 0.55,
    "Not every project needs eight stages. Templates — Standard Dev Flow (8) · Support Fast-Track (5) · Infra Flow (5) — "
    "are mapped per project × task type. The ribbon renders whatever the template defines; no code change to add a stage.",
    size=12.5, color=PRIMARY, line_spacing=1.3)
pg(s)

# ---------------------------------------------------------------- 11 loopbacks + counters
s = slide()
y = header(s, "Loop-backs, and two counters that are easy to confuse", "")
txt(s, ML, y, 5.9, 0.3, "LOOP-BACKS  ·  each increments the iteration counter", size=10.5, bold=True, color=TEXT2)
rows = [
    [("Rework", True, TEXT), "QA → Development", ("REWORK", False, TEXT2)],
    [("Deployment failure", True, TEXT), "Deployment → Development", ("DEPLOY_FAILED", False, TEXT2)],
    [("Verification failure", True, TEXT), "Verification → Development", ("VERIFY_FAILED", False, TEXT2)],
    [("Rejected at sign-off", True, TEXT), "Sign-off → Development", ("SIGNOFF_REJECTED", False, TEXT2)],
    [("Clarification", True, TEXT), "Development → Triage", ("CLARIFICATION", False, TEXT2)],
    [("Force move  (PM/Admin)", True, TEXT), "any → any", ("OVERRIDE", False, TEXT2)],
]
table(s, ML, y+0.36, 5.9, ["Loop", "From → To", "Code"], rows,
      [1.95, 2.3, 1.65], row_h=0.4, head_h=0.4, fsize=11)
txt(s, ML, y+3.15, 5.9, 0.3, "Every backward move requires a reason.", size=12, bold=True, color=DANGER)

x2 = ML + 6.2
txt(s, x2, y, CW-6.2, 0.3, "ITERATIONS vs CYCLES", size=10.5, bold=True, color=TEXT2)
card(s, x2, y+0.36, CW-6.2, 1.42, fill=RGBColor(0xFE,0xF3,0xC7), line=None)
txt(s, x2+0.3, y+0.56, CW-6.8, 0.3, "ITERATION  (iteration_no)", size=12, bold=True, color=RGBColor(0x92,0x40,0x0E))
txt(s, x2+0.3, y+0.92, CW-6.8, 0.75,
    "Increments when the ticket is pushed backwards in the ribbon.\nScope: within one cycle.",
    size=12.5, color=RGBColor(0x92,0x40,0x0E), line_spacing=1.28)
card(s, x2, y+1.95, CW-6.2, 1.42, fill=PRIMARY_SOFT, line=None)
txt(s, x2+0.3, y+2.15, CW-6.8, 0.3, "CYCLE  (cycle_no)", size=12, bold=True, color=PRIMARY)
txt(s, x2+0.3, y+2.51, CW-6.8, 0.75,
    "Increments when the ticket is reopened after closure.\nScope: the whole ticket. Each cycle has its own ribbon.",
    size=12.5, color=PRIMARY, line_spacing=1.28)
card(s, x2, y+3.55, CW-6.2, 0.92, fill=RGBColor(0x1E,0x1B,0x33), line=None)
txt(s, x2+0.3, y+3.75, CW-6.8, 0.6,
    "Cycle 2 · Iteration 3 · currently in QA\nClosed once, reopened, and QA has already bounced it twice this life.",
    size=12.5, color=WHITE, line_spacing=1.3)
pg(s)

# ---------------------------------------------------------------- 12 segment states
s = slide()
y = header(s, "What the ribbon renders", "Every segment shows five things · six visual states")
card(s, ML, y, 4.3, 2.35, fill=RGBColor(0x1E,0x1B,0x33), line=None)
txt(s, ML+0.35, y+0.25, 3.8, 1.9,
    "┌───────────────┐\n"
    "│ ✔ Development │  ← stage + state icon\n"
    "│ Ravi Kumar    │  ← owner\n"
    "│ 3d 4h         │  ← time in stage\n"
    "│ 14.5 h effort │  ← effort in stage\n"
    "│ ↺ ×2          │  ← loop-back badge\n"
    "└───────────────┘",
    size=11.5, color=RGBColor(0xE5,0xE8,0xF0), font=MONO, line_spacing=1.28)
rows = [
    [("Completed", True, SUCCESS), "Solid tick, green accent, filled connector"],
    [("Current", True, PRIMARY), "Indigo fill, pulse ring, \"Now\" label, live elapsed timer"],
    [("Pending", True, TEXT2), "Outlined, muted text, dashed connector"],
    [("Reworked", True, WARNING), "Amber left edge + ↺ ×N badge"],
    [("Skipped", True, TEXT2), "Dashed outline, strikethrough, hover shows who authorised"],
    [("Blocked / breached", True, DANGER), "Grey with pause icon, or red past stage SLA"],
]
table(s, ML+4.6, y, CW-4.6, ["Segment state", "Visual treatment"], rows,
      [2.4, 4.93], row_h=0.36, head_h=0.4, fsize=11.5)
card(s, ML, y+2.65, CW, 1.85, fill=SURFACE)
txt(s, ML+0.35, y+2.85, CW-0.7, 0.3, "INTERACTIONS", size=11, bold=True, color=TEXT2)
bullets(s, ML+0.35, y+3.2, CW-0.8, [
    ("Click a segment ", "— filters the History, Effort and Chat tabs below to just that stage and iteration."),
    ("Current segment ", "— carries the contextual action inline: Hand off to QA → · Pass / Fail · Mark deployed. Hidden for everyone else."),
    ("Compact ribbon in lists ", "— eight dots per row, so a manager scans a whole grid and sees where every ticket sits."),
], size=12, gap=0.4)
pg(s)

# ---------------------------------------------------------------- 13 effort attribution
s = slide()
y = header(s, "Effort attribution per resource, per stage",
           "The panel under the ribbon — \"what did your team actually do on this issue?\"")
rows = [
    [("1", False, TEXT2), "Intake", "Priya N.", ("SUP", False, TEXT2), "1h 10m", ("0.5 h", True, TEXT), ""],
    [("1", False, TEXT2), "Triage", "Meera P.", ("PM", False, TEXT2), "3h 20m", ("1.0 h", True, TEXT), ""],
    [("1", False, TEXT2), "Development", "Ravi K.", ("DEV", False, TEXT2), "2d 1h", ("9.0 h", True, TEXT), ""],
    [("1", False, TEXT2), "QA", "Anil S.", ("QA", False, TEXT2), "6h 40m", ("3.5 h", True, TEXT), ("✗ Failed", True, DANGER)],
    [("2", False, TEXT2), "Development", "Ravi K.", ("DEV", False, TEXT2), "1d 2h", ("5.5 h", True, TEXT), ("↺ rework", True, WARNING)],
    [("2", False, TEXT2), "QA", "Anil S.", ("QA", False, TEXT2), "4h 05m", ("2.0 h", True, TEXT), ("✓ Passed", True, SUCCESS)],
    [("2", False, TEXT2), "Deployment", "Karan D.", ("DEP", False, TEXT2), "5h 30m", ("1.5 h", True, TEXT), ""],
    [("", False, TEXT2), ("TOTAL — Cycle 2", True, TEXT), "", "", ("5d 3h", True, TEXT), ("24.5 h", True, PRIMARY), ""],
    [("", False, TEXT2), ("TOTAL — all cycles", True, TEXT), "", "", ("12d 6h", True, TEXT), ("38.0 h", True, PRIMARY), ""],
]
table(s, ML, y, CW, ["It", "Stage", "Resource", "Role", "Duration", "Effort", ""], rows,
      [0.5, 2.4, 2.0, 1.0, 1.9, 1.6, 2.53], row_h=0.32, head_h=0.38, fsize=11)
card(s, ML, y+3.5, CW, 1.0, fill=RGBColor(0xEC,0xFD,0xF5), line=None)
txt(s, ML+0.35, y+3.68, CW-0.7, 0.65,
    "Active time = Σ effort logged in that stage.    Idle / queue time = duration in stage − active time.\n"
    "A stage with 2 days duration and 2 hours of effort is a QUEUE problem, not a capacity problem — and that "
    "single insight is usually worth the whole ribbon feature.",
    size=12.5, color=RGBColor(0x06,0x5F,0x46), line_spacing=1.32)
pg(s)

# ---------------------------------------------------------------- 14 ticket page
s = slide()
y = header(s, "The ticket page — five capabilities that make it usable",
           "Individually small; together they are what a support desk handling real client traffic needs")
items = [
    ("Priority dropdown", "Colour chips, not plain text. Pre-filled from the task type. Changing it recomputes and previews the planned close date before you commit. Mandatory reason once assigned.", A_C),
    ("Client + contact", "Type-ahead over name, code and domain, filtered to the project's clients. Dependent contact dropdown with inline \"+ Add contact\" so the desk never leaves the form.", B_C),
    ("Attachments", "Drag-drop, file picker, and paste from clipboard — the last matters most, because pasting a screenshot straight from Snipping Tool is the most common attachment action.", C_C),
    ("Comment box", "Rich text, @mentions, internal vs client-visible toggle, 5-minute edit window then locked. Stamped with the stage and iteration it was written in.", D_C),
    ("Mail alert engine", "A mail fires on every stage transition and every event that changes who is responsible. Queued, templated, threaded, logged, retried.", A_C),
]
cy = y
for t, b, col in items:
    card(s, ML, cy, CW, 0.83)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML), Inches(cy+0.1), Inches(0.055), Inches(0.63))
    bar.fill.solid(); bar.fill.fore_color.rgb = col
    bar.line.fill.background(); bar.shadow.inherit = False
    txt(s, ML+0.3, cy+0.28, 2.5, 0.3, t, size=14, bold=True, color=col)
    txt(s, ML+3.0, cy+0.16, CW-3.3, 0.62, b, size=11.5, color=TEXT2, line_spacing=1.24)
    cy += 0.93
pg(s)

# ---------------------------------------------------------------- 15 comments history
s = slide()
y = header(s, "Comments interleave into one chronological history",
           "So \"what happened on this ticket\" is a single readable timeline, not two lists to reconcile")
code_block(s, ML, y, CW, 2.75,
           "06 Aug 14:22   💬  Ravi Kumar (Developer) · Development · iteration 2\n"
           "               \"Root cause is the retry timeout. Patch pushed, ready for QA.\"       [internal]\n"
           "\n"
           "06 Aug 14:25   ➡  Handoff   Development → QA    Ravi Kumar → Anil Sharma\n"
           "               Note: \"Please retest the checkout path on staging\"\n"
           "\n"
           "06 Aug 16:40   🔺  Level changed   High → Critical    by System (SLA breach)\n"
           "\n"
           "07 Aug 09:15   💬  Anil Sharma (QA) · QA · iteration 2\n"
           "               \"Verified on staging, three of three defects closed.\"                [internal]\n"
           "\n"
           "07 Aug 09:16   📎  Attachment added   qa-signoff-report.pdf (412 KB)   by Anil Sharma",
           size=11)
card(s, ML, y+3.05, 5.85, 1.42, fill=PRIMARY_SOFT, line=None)
txt(s, ML+0.3, y+3.25, 5.25, 1.05,
    "Comments are the permanent record.\n\n"
    "Chat is conversational and ephemeral; comments are auditable, stamped with stage and iteration, "
    "and immutable after five minutes.",
    size=12.5, color=PRIMARY, line_spacing=1.3)
card(s, ML+6.15, y+3.05, CW-6.15, 1.42, fill=RGBColor(0xFE,0xF2,0xF2), line=None)
txt(s, ML+6.45, y+3.25, CW-6.75, 1.05,
    "Default to internal, always.\n\n"
    "An accidental leak to a client is far costlier than an extra click. "
    "Client-visible renders in a different colour before you post.",
    size=12.5, color=RGBColor(0x99,0x1B,0x1B), line_spacing=1.3)
pg(s)

# ---------------------------------------------------------------- 16 client + import
s = slide()
y = header(s, "Client master & Excel bulk import",
           "A five-step wizard, because a silent import that half-succeeds is worse than no import at all")
steps = [
    ("1", "DOWNLOAD TEMPLATE", "Pre-formatted .xlsx with exact headers, validation dropdowns and one example row"),
    ("2", "UPLOAD", "Drag-drop .xlsx/.csv, max 5 MB, up to 5,000 rows, parsed server-side"),
    ("3", "MAP COLUMNS", "Auto-matched by header with manual override. Presets saved for next time"),
    ("4", "VALIDATE + PREVIEW", "Dry run — nothing written. Per-row: will create · will update · duplicate · rejected + reason"),
    ("5", "COMMIT", "Background job with progress. Downloadable error report for every rejected row"),
]
cy = y
for num, t, b in steps:
    card(s, ML, cy, 6.6, 0.6)
    nn = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(ML+0.22), Inches(cy+0.15), Inches(0.32), Inches(0.32))
    nn.fill.solid(); nn.fill.fore_color.rgb = PRIMARY
    nn.line.fill.background(); nn.shadow.inherit = False
    txt(s, ML+0.22, cy+0.21, 0.32, 0.26, num, size=11.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, ML+0.68, cy+0.09, 2.2, 0.28, t, size=11.5, bold=True, color=TEXT)
    txt(s, ML+0.68, cy+0.33, 5.6, 0.24, b, size=10, color=TEXT2)
    cy += 0.68
x2 = ML + 6.9
card(s, x2, y, CW-6.9, 1.75, fill=RGBColor(0xEC,0xFD,0xF5), line=None)
txt(s, x2+0.3, y+0.2, CW-7.5, 0.3, "VALIDATION RULES", size=10.5, bold=True, color=RGBColor(0x06,0x5F,0x46))
txt(s, x2+0.3, y+0.58, CW-7.5, 1.1,
    "Client code unique · name required · email format · country against ISO list · "
    "duplicate detection on code then email domain · existing records UPDATED, never duplicated (upsert on code)",
    size=11.5, color=RGBColor(0x06,0x5F,0x46), line_spacing=1.28)
card(s, x2, y+1.95, CW-6.9, 1.45, fill=PRIMARY_SOFT, line=None)
txt(s, x2+0.3, y+2.15, CW-7.5, 1.1,
    "Every import writes an import_batch row — so a bad import can be identified and reversed as a set.\n\n"
    "Same wizard registered twice: clients and resources.",
    size=11.5, color=PRIMARY, line_spacing=1.3)
card(s, x2, y+3.6, CW-6.9, 0.87, fill=SURFACE)
txt(s, x2+0.3, y+3.78, CW-7.5, 0.55,
    "Deactivating a client with open tickets warns and blocks new tickets — but never hides historical ones.",
    size=11.5, color=TEXT, line_spacing=1.28)
pg(s)

# ---------------------------------------------------------------- 17 mail engine
s = slide()
y = header(s, "Mail is the guaranteed channel",
           "In-app popups only reach someone who is logged in")
rows = [
    [("Ticket assigned to you", True, TEXT), "Assignee", ("[CRM-26-00347] New ticket assigned — Critical", False, TEXT2), ("never optional", True, DANGER)],
    [("Handoff — ribbon moves to you", True, TEXT), "New stage owner", ("[CRM-26-00347] Handed to you at QA by Ravi Kumar", False, TEXT2), ("never optional", True, DANGER)],
    [("Sent back for rework", True, TEXT), "Developer, cc PM", ("[CRM-26-00347] QA failed — 3 defects returned", False, TEXT2), ("never optional", True, DANGER)],
    [("Escalated to Critical", True, TEXT), "Assignee, PM, RM", ("[CRM-26-00347] Escalated to CRITICAL", False, TEXT2), ("never optional", True, DANGER)],
    [("SLA breached / delayed", True, TEXT), "Assignee, PM, RM", ("[CRM-26-00347] Overdue by 3 days", False, TEXT2), ("never optional", True, DANGER)],
    [("Comment added · @mention", True, TEXT), "Assignee, watchers", ("[CRM-26-00347] New comment from Meera Prasad", False, TEXT2), ("digest option", True, SUCCESS)],
    [("Daily digest 08:30", True, TEXT), "Everyone", ("Your open tickets — 4 due today, 1 overdue", False, TEXT2), ("opt-out", True, SUCCESS)],
]
table(s, ML, y, CW, ["Event", "To", "Subject pattern", "Optional?"], rows,
      [3.3, 2.2, 4.7, 1.73], row_h=0.38, head_h=0.4, fsize=11)
card(s, ML, y+3.3, CW, 1.2, fill=RGBColor(0x1E,0x1B,0x33), line=None)
txt(s, ML+0.35, y+3.48, CW-0.7, 0.3, "ENGINEERING RULES", size=10.5, bold=True, color=RGBColor(0xA5,0xB4,0xFC))
txt(s, ML+0.35, y+3.83, CW-0.7, 0.6,
    "Queued, never sent inline — a slow SMTP server must never slow a handoff.   ·   Message-ID / In-Reply-To keyed on the "
    "ticket, so a whole ticket collapses into one thread.   ·   Every send logged; 3 retries; bounce webhooks.   ·   "
    "One mail per recipient per ticket per minute.",
    size=12, color=WHITE, line_spacing=1.3)
pg(s)

# ---------------------------------------------------------------- 18 formulas
s = slide()
y = header(s, "Effort, velocity and SLA — the calculations", "")
code_block(s, ML, y, 5.85, 3.15,
           "Effort variance %  = (logged − estimated) / estimated × 100\n"
           "Cycle time         = actual_close − cycle.start   (working hrs)\n"
           "Lead time          = actual_close − date_reported\n"
           "SLA compliance %   = closed_within_pcd / total_closed × 100\n"
           "Reopen rate %      = tickets_reopened / total_closed × 100\n"
           "\n"
           "Time in stage      = exited_at − entered_at\n"
           "Active in stage    = Σ effort with that stage + iteration\n"
           "Idle / queue time  = time in stage − active time\n"
           "Rework rate %      = iteration_no > 1 / total × 100\n"
           "Handoff latency    = next entered_at − prev exited_at\n"
           "First-time-right % = closed at iteration 1 / total closed", size=11)
x2 = ML + 6.15
card(s, x2, y, CW-6.15, 1.5, fill=PRIMARY_SOFT, line=None)
txt(s, x2+0.3, y+0.22, CW-6.75, 0.3, "RESOURCE VELOCITY", size=10.5, bold=True, color=PRIMARY)
txt(s, x2+0.3, y+0.6, CW-6.75, 0.8,
    "Tickets closed per week (count) AND effort-weighted: Σ effort-hours closed per week.\n"
    "Utilisation % = logged hours / available hours.",
    size=12, color=PRIMARY, line_spacing=1.3)
card(s, x2, y+1.7, CW-6.15, 1.45, fill=RGBColor(0xFE,0xF3,0xC7), line=None)
txt(s, x2+0.3, y+1.92, CW-6.75, 0.3, "THE WORKING CALENDAR", size=10.5, bold=True, color=RGBColor(0x92,0x40,0x0E))
txt(s, x2+0.3, y+2.3, CW-6.75, 0.8,
    "Weekends + org holidays + resource leave — so a ticket raised Friday 6 PM with a 4-hour SLA isn't "
    "breached on Saturday morning.",
    size=12, color=RGBColor(0x92,0x40,0x0E), line_spacing=1.3)
card(s, ML, y+3.35, CW, 1.1, fill=SURFACE)
txt(s, ML+0.35, y+3.53, CW-0.7, 0.3, "ESCALATION ENGINE  ·  scheduled worker every 15 minutes", size=11, bold=True, color=TEXT2)
txt(s, ML+0.35, y+3.88, CW-0.7, 0.5,
    "80% of SLA elapsed → warning   ·   past planned close → level becomes CRITICAL, alert to Reporting Manager   ·   "
    "delayed > 48 h → escalate to L2   ·   no update in 3 working days → nudge   ·   unassigned > 2 h → triage alert",
    size=12, color=TEXT, line_spacing=1.3)
pg(s)

# ---------------------------------------------------------------- 19 screens
s = slide()
y = header(s, "34 screens across 8 modules",
           "Each maps 1:1 to a Figma frame and a Jira epic, so design, dev and QA share one identifier")
mods = [
    ("Authentication", "S-01 … S-04", "Login · forgot/reset · change password · 2FA", A_C),
    ("Common shell", "—", "Sidebar · top bar · global search · notification bell · toast layer", B_C),
    ("Dashboard", "S-05, S-06", "Role-aware dashboard · chart drill-down modal", C_C),
    ("Master data", "S-07 … S-16, S-32 … S-34", "Resource · role · project · task type · priority · client · import wizard · workflow templates · calendar · notification templates · audit log", D_C),
    ("Tickets", "S-17 … S-24, S-29 … S-31", "List · My Tasks · create · detail · quick update · reopen · close · bulk reassign · handoff · template designer · stage queue", A_C),
    ("Chat", "S-25", "Ticket thread · direct message · project channel · Ask Status", B_C),
    ("Notifications", "S-26", "Notification centre · per-user preference matrix", C_C),
    ("Reports", "S-27, S-28", "Reports hub — 18 reports · Resource 360° profile", D_C),
]
cy = y
for name, ids, body, col in mods:
    card(s, ML, cy, CW, 0.53)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML), Inches(cy+0.08), Inches(0.05), Inches(0.37))
    bar.fill.solid(); bar.fill.fore_color.rgb = col
    bar.line.fill.background(); bar.shadow.inherit = False
    txt(s, ML+0.26, cy+0.15, 1.9, 0.28, name, size=12, bold=True, color=TEXT)
    txt(s, ML+2.25, cy+0.15, 2.3, 0.28, ids, size=10.5, color=col, font=MONO)
    txt(s, ML+4.7, cy+0.13, CW-5.0, 0.36, body, size=10.5, color=TEXT2, line_spacing=1.15)
    cy += 0.6
pg(s)

# ---------------------------------------------------------------- 20 dashboard
s = slide()
y = header(s, "Dashboard — 20 widgets, every one drill-down enabled",
           "Each card and chart segment deep-links to a pre-filtered ticket list")
groups = [
    ("KPI cards", "Total created · Open/pending · Closed · Critical · Delayed · Reopened", A_C),
    ("Distribution", "Task-type donut · Priority bar · Project treemap · Aging buckets", B_C),
    ("Trend", "Daily status stacked area · Resource velocity multi-line · Date-wise heatmap · Handoff latency", C_C),
    ("Load & SLA", "Resource-wise load bar · SLA compliance radial gauge · Client-wise volume", D_C),
    ("Ribbon-driven", "Stage funnel (WIP) · Rework / ping-pong tickets · Avg time per stage, active vs idle", PRIMARY),
]
cy = y
for t, b, col in groups:
    card(s, ML, cy, 7.3, 0.66)
    txt(s, ML+0.3, cy+0.11, 2.2, 0.28, t, size=12.5, bold=True, color=col)
    txt(s, ML+0.3, cy+0.38, 6.7, 0.24, b, size=10.5, color=TEXT2)
    cy += 0.74
x2 = ML + 7.6
card(s, x2, y, CW-7.6, 1.8, fill=PRIMARY_SOFT, line=None)
txt(s, x2+0.3, y+0.22, CW-8.2, 1.4,
    "Developer's dashboard shows only 8 of the 20, scoped to assignee = me — plus \"My due today / this week\".\n\n"
    "Role-awareness is enforced server-side, not by hiding cards.",
    size=12, color=PRIMARY, line_spacing=1.3)
card(s, x2, y+2.0, CW-7.6, 1.7, fill=RGBColor(0xEC,0xFD,0xF5), line=None)
txt(s, x2+0.3, y+2.2, CW-8.2, 1.35,
    "Reads come from pre-aggregated summary tables refreshed every 5 minutes.\n\n"
    "Never a live COUNT(*) over the full ticket table.",
    size=12, color=RGBColor(0x06,0x5F,0x46), line_spacing=1.3)
pg(s)

# ---------------------------------------------------------------- 21 reports
s = slide()
y = header(s, "18 reports — all filterable, exportable and schedulable", "Excel · CSV · PDF, with daily / weekly / monthly email scheduling")
rows = [
    [("Resource Performance Scorecard", True, TEXT), "Assigned, closed, on-time, SLA %, cycle time, variance, reopen rate, utilisation"],
    [("Resource Velocity", True, TEXT), "Tickets & effort-hours closed per week, 4-week rolling average, multi-resource comparison"],
    [("Effort Summary", True, TEXT), "Effort by resource × project × task type, pivot-style with drill-down to individual logs"],
    [("Delayed / SLA Breach", True, TEXT), "Every breach, days overdue, escalation level, reason"],
    [("Stage Funnel / WIP · Stage Cycle Time", True, TEXT), "Live count per stage, bottleneck highlighted; avg time split into active vs idle"],
    [("Rework Analysis", True, TEXT), "Rework rate by developer, QA rejection rate, first-time-right %"],
    [("Resource Contribution", True, TEXT), "The per-resource-per-stage roll-up, across any ticket set"],
    [("Audit / Compliance", True, TEXT), "Full immutable trail export for a ticket or date range, including every handoff"],
    [("Client Report · Email Delivery Log", True, TEXT), "Volume, SLA and satisfaction per client; every alert sent, with delivery status"],
]
table(s, ML, y, CW, ["Report", "What it answers"], rows,
      [4.2, 7.73], row_h=0.36, head_h=0.4, fsize=11)
card(s, ML, y+3.66, CW, 0.82, fill=PRIMARY_SOFT, line=None)
txt(s, ML+0.35, y+3.84, CW-0.7, 0.52,
    "Plus: Task Type Analysis · Reopen Analysis · Date-wise · Project Health · Aging · Workload/Capacity · Deployment Report.   "
    "Resource 360° Profile reaches all of it in one click from anywhere a name appears.",
    size=12, color=PRIMARY, line_spacing=1.3)
pg(s)

# ---------------------------------------------------------------- 22 design system
s = slide()
y = header(s, "Design system — light theme only", "WCAG AA on every token, keyboard navigable, prefers-reduced-motion respected")
sw = [("#F7F8FC","bg-app"),("#FFFFFF","surface"),("#F1F3F9","subtle"),("#E5E8F0","border"),
      ("#4F46E5","primary"),("#EEF2FF","primary-soft"),("#111827","text"),("#6B7280","text-2"),
      ("#10B981","success"),("#F59E0B","warning"),("#EF4444","danger"),("#3B82F6","info")]
cy = y
for i, (hexv, name) in enumerate(sw):
    x = ML + (i % 6) * 2.0
    yy = cy + (i // 6) * 1.05
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(yy), Inches(1.8), Inches(0.6))
    c.fill.solid(); c.fill.fore_color.rgb = RGBColor.from_string(hexv[1:])
    c.line.color.rgb = BORDER; c.line.width = Pt(1); c.shadow.inherit = False
    txt(s, x, yy+0.66, 1.8, 0.22, name, size=10.5, bold=True, color=TEXT)
    txt(s, x, yy+0.86, 1.8, 0.2, hexv, size=9.5, color=TEXT2, font=MONO)
card(s, ML, y+2.35, 5.85, 2.1, fill=SURFACE)
txt(s, ML+0.3, y+2.55, 5.25, 0.3, "FOUNDATIONS", size=10.5, bold=True, color=TEXT2)
txt(s, ML+0.3, y+2.92, 5.25, 1.4,
    "Type   Inter / Plus Jakarta Sans · base 14/20\n"
    "Space  4px scale · card padding 20 · gutter 24\n"
    "Radius 12 cards · 8 inputs · 999 chips\n"
    "Motion 150–200 ms ease-out · slide-over from right\n"
    "Loading  skeletons, never spinners",
    size=11.5, color=TEXT, line_spacing=1.35)
card(s, ML+6.15, y+2.35, CW-6.15, 2.1, fill=RGBColor(0x1E,0x1B,0x33), line=None)
txt(s, ML+6.45, y+2.55, CW-6.75, 0.3, "NEVER COLOUR ALONE", size=10.5, bold=True, color=RGBColor(0xA5,0xB4,0xFC))
txt(s, ML+6.45, y+2.95, CW-6.75, 1.35,
    "Every ribbon state carries an icon and a text label as well as a colour — so it stays readable for "
    "colour-blind users and in print.\n\n"
    "Chart palette is colour-blind safe.",
    size=12.5, color=WHITE, line_spacing=1.32)
pg(s)

# ---------------------------------------------------------------- 23 roadmap
s = slide()
y = header(s, "Delivery roadmap", "Roughly five months to go-live with a seven-person team")
rows = [
    [("0", True, PRIMARY), ("Foundation", True, TEXT), ("1–2", False, TEXT2), "Repo, CI/CD, DB schema + migrations, design system, auth skeleton"],
    [("1", True, PRIMARY), ("Core MVP", True, TEXT), ("3–7", False, TEXT2), "Login, masters, client master + Excel import, ticket CRUD with priority/client/attachments, comments, cycle + history + effort, workflow templates, ribbon, Journey tab"],
    [("2", True, PRIMARY), ("Intelligence", True, TEXT), ("8–11", False, TEXT2), "Dashboard + drill-down, SLA engine, stage SLA alerts, escalation, notification centre, full mail engine, stage queue"],
    [("3", True, PRIMARY), ("Collaboration", True, TEXT), ("12–14", False, TEXT2), "Chat (ticket / DM / project), Ask Status, mentions, watchers, attachments"],
    [("4", True, PRIMARY), ("Reporting", True, TEXT), ("15–17", False, TEXT2), "Reports hub, scorecard, velocity, stage cycle time, rework analysis, exports, Resource 360"],
    [("5", True, PRIMARY), ("Hardening", True, TEXT), ("18–20", False, TEXT2), "Audit viewer, hash-chain verification, performance, security audit, UAT, training, go-live"],
    [("6", True, TEXT2), ("Extensions", True, TEXT), ("post", False, TEXT2), "Client portal, email-to-ticket, knowledge base, mobile app, SSO, Teams/Slack"],
]
table(s, ML, y, CW, ["", "Phase", "Weeks", "Scope"], rows,
      [0.5, 1.9, 1.0, 8.53], row_h=0.5, head_h=0.4, fsize=11)
card(s, ML, y+4.0, CW, 0.6, fill=PRIMARY_SOFT, line=None)
txt(s, ML+0.35, y+4.14, CW-0.7, 0.35,
    "Team: 1 PM/BA · 1 architect · 2 backend · 2 frontend · 1 QA · 0.5 DevOps · 0.5 designer.   "
    "The ribbon adds ~2–3 weeks of the phase-1 budget.",
    size=12, color=PRIMARY)
pg(s)

# ---------------------------------------------------------------- 24 risks
s = slide()
y = header(s, "Risks and mitigations", "")
rows = [
    [("History integrity challenged in a client dispute", True, TEXT), "Hash-chained append-only tables, nightly verification, audit export"],
    [("Notification fatigue → people ignore alerts", True, TEXT), "Per-user preference matrix, digests, escalations reserved for genuine breaches"],
    [("Dashboard slows as data grows", True, TEXT), "Pre-aggregated summary tables, partial indexes, Redis cache, archival after 3 years"],
    [("Resources under-log effort → velocity is fiction", True, TEXT), "Effort mandatory at handoff, daily reminders, manager approval"],
    [("Role scope bug leaks another team's tickets", True, TEXT), "One central guard + automated permission tests covering every role × route"],
    [("Ribbon unreadable at 8 stages on a laptop", True, TEXT), "Compact dot variant, horizontal scroll with current segment auto-centred, collapsed completed stages"],
    [("Teams game the ribbon to stop their stage clock", True, TEXT), "Idle-vs-active reported alongside duration; iteration counter makes ping-pong visible within a day"],
    [("Assignee never sees the alert", True, TEXT), "Mail + in-app + queued popup on next login, with a per-ticket delivery log — a missed alert is provable, not deniable"],
    [("Internal debug notes leak to a client", True, TEXT), "Comments and attachments default to internal; client-visible is an explicit, differently-coloured toggle"],
]
table(s, ML, y, CW, ["Risk", "Mitigation"], rows, [5.2, 6.73], row_h=0.36, head_h=0.4, fsize=11)
pg(s)

# ---------------------------------------------------------------- 25 governance
s = slide(SURFACE)
b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SW), Inches(SH))
b.fill.gradient(); b.fill.gradient_angle = 45.0
st = b.fill.gradient_stops
st[0].color.rgb = RGBColor(0x4F,0x46,0xE5); st[0].position = 0.0
st[1].color.rgb = RGBColor(0x1E,0x1B,0x4B); st[1].position = 1.0
b.line.fill.background(); b.shadow.inherit = False
txt(s, ML+0.3, 0.75, CW, 0.5, "Governance decisions to lock before build", size=30, bold=True, color=WHITE)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML+0.3), Inches(1.42), Inches(1.3), Inches(0.04))
ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0xA5,0xB4,0xFC)
ln.line.fill.background(); ln.shadow.inherit = False
qs = [
    ("Can a Developer close a ticket?", "Resolved only — closure belongs to the Sign-off stage owner"),
    ("Can a ticket skip QA?", "PM/Admin only, reason mandatory, never for Production Bug"),
    ("Does rework reset the Planned Close Date?", "No — the original date stands; that is the honest measure"),
    ("Is effort mandatory at handoff?", "Yes, blocking — otherwise the roll-up degrades within weeks"),
    ("Comments default internal or client-visible?", "Internal, always — a leak costs more than a click"),
    ("Does an import update or only insert?", "Upsert on client code, with the dry run showing what changes"),
    ("Does escalated Critical revert after closure?", "Keep original_level and report both"),
    ("Is effort self-reported or timer-based?", "Self-reported with optional timer — timers get forgotten"),
]
cy = 1.85
for q, a in qs:
    txt(s, ML+0.3, cy, 5.6, 0.3, q, size=13, bold=True, color=WHITE)
    txt(s, ML+6.1, cy+0.02, CW-6.4, 0.3, "→  " + a, size=12.5, color=RGBColor(0xC7,0xD2,0xFE))
    cy += 0.6
txt(s, ML+0.3, 6.85, CW, 0.3,
    "Full list in blueprint §16 · build sequencing in PLAN.md · stream assignment in TEAM-PLAN.md",
    size=11.5, color=RGBColor(0x81,0x8C,0xF8))

save(str(pathlib.Path(__file__).parents[2] / "docs/decks") + "/EduTrack-Product-Blueprint.pptx")
