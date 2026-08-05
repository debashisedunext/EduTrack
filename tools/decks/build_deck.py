#!/usr/bin/env python3
"""Generate EduTrack-Team-Plan.pptx from TEAM-PLAN.md content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------- tokens
# from blueprint 12.1
PRIMARY      = RGBColor(0x4F, 0x46, 0xE5)
PRIMARY_SOFT = RGBColor(0xEE, 0xF2, 0xFF)
BG_APP       = RGBColor(0xF7, 0xF8, 0xFC)
SURFACE      = RGBColor(0xFF, 0xFF, 0xFF)
SUBTLE       = RGBColor(0xF1, 0xF3, 0xF9)
BORDER       = RGBColor(0xE5, 0xE8, 0xF0)
TEXT         = RGBColor(0x11, 0x18, 0x27)
TEXT2        = RGBColor(0x6B, 0x72, 0x80)
SUCCESS      = RGBColor(0x10, 0xB9, 0x81)
WARNING      = RGBColor(0xF5, 0x9E, 0x0B)
DANGER       = RGBColor(0xEF, 0x44, 0x44)
INFO         = RGBColor(0x3B, 0x82, 0xF6)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)

# stream accents
A_C = RGBColor(0x4F, 0x46, 0xE5)
B_C = RGBColor(0x06, 0xB6, 0xD4)
C_C = RGBColor(0x8B, 0x5C, 0xF6)
D_C = RGBColor(0x14, 0xB8, 0xA6)

FONT = "Calibri"
MONO = "Consolas"

SW, SH = 13.333, 7.5
ML = 0.7                      # left margin
CW = SW - 2 * ML              # content width

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def slide(bg=BG_APP):
    s = prs.slides.add_slide(BLANK)
    bgshape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                 Inches(SW), Inches(SH))
    bgshape.fill.solid()
    bgshape.fill.fore_color.rgb = bg
    bgshape.line.fill.background()
    bgshape.shadow.inherit = False
    return s


def txt(s, x, y, w, h, text, size=14, bold=False, color=TEXT,
        align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP, space_after=0,
        line_spacing=1.15):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def header(s, title, sub=None, accent=PRIMARY):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML), Inches(0.52),
                             Inches(0.055), Inches(0.46))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False
    txt(s, ML + 0.2, 0.5, CW - 0.2, 0.5, title, size=27, bold=True, color=TEXT)
    if sub:
        txt(s, ML + 0.2, 1.05, CW - 0.2, 0.36, sub, size=13.5, color=TEXT2)
        return 1.62
    return 1.28


def card(s, x, y, w, h, fill=SURFACE, line=BORDER, radius=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    if radius:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    return shp


def table(s, x, y, w, cols, rows, col_w, header_fill=PRIMARY,
          header_color=WHITE, fsize=11.5, hsize=11.5, row_h=0.34,
          head_h=0.4, zebra=True):
    """rows: list of list of (text, bold, color) or plain str"""
    n_rows = len(rows) + 1
    shape = s.shapes.add_table(n_rows, len(cols), Inches(x), Inches(y),
                               Inches(w), Inches(head_h + row_h * len(rows)))
    tbl = shape.table
    tbl.first_row = False
    tbl.horz_banding = False
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = Inches(cw)
    tbl.rows[0].height = Inches(head_h)
    for r in range(1, n_rows):
        tbl.rows[r].height = Inches(row_h)

    def fill_cell(cell, value, bold, color, fill_rgb, size):
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_rgb
        cell.margin_left = Inches(0.09)
        cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = 1.0
        r = p.add_run()
        r.text = value
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = FONT

    for c, head in enumerate(cols):
        fill_cell(tbl.cell(0, c), head, True, header_color, header_fill, hsize)
    for ri, row in enumerate(rows):
        bg = SURFACE if (not zebra or ri % 2 == 0) else SUBTLE
        for ci, val in enumerate(row):
            if isinstance(val, tuple):
                t, b, col = val
            else:
                t, b, col = val, False, TEXT
            fill_cell(tbl.cell(ri + 1, ci), t, b, col, bg, fsize)
    return shape


def bullets(s, x, y, w, items, size=14, gap=0.42, dot=PRIMARY, bold_lead=True):
    """items: list of (lead, rest) or plain string"""
    cy = y
    for it in items:
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(cy + 0.075),
                               Inches(0.085), Inches(0.085))
        d.fill.solid()
        d.fill.fore_color.rgb = dot
        d.line.fill.background()
        d.shadow.inherit = False
        tb = s.shapes.add_textbox(Inches(x + 0.22), Inches(cy),
                                  Inches(w - 0.22), Inches(gap))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.line_spacing = 1.2
        if isinstance(it, tuple):
            lead, rest = it
            r1 = p.add_run(); r1.text = lead
            r1.font.size = Pt(size); r1.font.bold = bold_lead
            r1.font.color.rgb = TEXT; r1.font.name = FONT
            r2 = p.add_run(); r2.text = rest
            r2.font.size = Pt(size); r2.font.color.rgb = TEXT2
            r2.font.name = FONT
        else:
            r = p.add_run(); r.text = it
            r.font.size = Pt(size); r.font.color.rgb = TEXT
            r.font.name = FONT
        cy += gap
    return cy


def footer(s, n, label="EduTrack — Team Plan"):
    txt(s, ML, SH - 0.46, 4.0, 0.26, label, size=9.5, color=TEXT2)
    txt(s, SW - ML - 1.2, SH - 0.46, 1.2, 0.26, str(n), size=9.5,
        color=TEXT2, align=PP_ALIGN.RIGHT)


def code_block(s, x, y, w, h, lines, size=11.5):
    card(s, x, y, w, h, fill=RGBColor(0x1E, 0x1B, 0x33), line=None)
    txt(s, x + 0.22, y + 0.16, w - 0.44, h - 0.32, lines, size=size,
        color=RGBColor(0xE5, 0xE8, 0xF0), font=MONO, line_spacing=1.28)


# ================================================================ 1 title
s = slide(SURFACE)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SW), Inches(SH))
band.fill.gradient()
band.fill.gradient_angle = 45.0
stops = band.fill.gradient_stops
stops[0].color.rgb = RGBColor(0x4F, 0x46, 0xE5)
stops[0].position = 0.0
stops[1].color.rgb = RGBColor(0x1E, 0x1B, 0x4B)
stops[1].position = 1.0
band.line.fill.background()
band.shadow.inherit = False

txt(s, ML + 0.3, 2.05, CW, 0.4, "EDUTRACK", size=15, bold=True,
    color=RGBColor(0xA5, 0xB4, 0xFC))
txt(s, ML + 0.3, 2.55, CW, 1.5,
    "Team Structure, Work Division\n& Git Workflow",
    size=42, bold=True, color=WHITE, line_spacing=1.08)
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML + 0.3), Inches(4.35),
                          Inches(1.5), Inches(0.045))
line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0xA5, 0xB4, 0xFC)
line.line.fill.background(); line.shadow.inherit = False
txt(s, ML + 0.3, 4.68, CW - 1.0, 0.9,
    "Four developers · four vertical streams · 18 weeks\n"
    "Java 21 · Spring Boot · MySQL 8.4 · React 18",
    size=15.5, color=RGBColor(0xC7, 0xD2, 0xFE), line_spacing=1.35)
txt(s, ML + 0.3, 6.5, CW, 0.3, "Companion to PLAN.md and the Product Blueprint v1.2",
    size=11.5, color=RGBColor(0x81, 0x8C, 0xF8))

# ================================================================ 2 problem
s = slide()
y = header(s, "The problem this plan solves",
           "Four developers on one codebase produce three predictable failures")
cx = ML
cwid = (CW - 0.5) / 3
probs = [
    ("01", "The blocking chain", DANGER,
     "PLAN.md sequences M0 → M1 → M2 → M3/M4 by dependency. Read literally, "
     "three developers sit idle for six weeks waiting on the schema and the "
     "security spine.", "Sprint 0 + three decouplers"),
    ("02", "Merge collisions", WARNING,
     "A conventional Spring layout puts every controller in one folder. Four "
     "developers then edit the same four directories every single day.",
     "Feature packaging + ownership map"),
    ("03", "Contract drift", INFO,
     "Four people building against an API that does not exist yet will each "
     "assume a different shape for it, and discover the mismatch in week eight.",
     "OpenAPI spec as a Sprint 0 deliverable"),
]
for i, (num, title, col, body, fix) in enumerate(probs):
    x = cx + i * (cwid + 0.25)
    card(s, x, y, cwid, 4.15)
    tag = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.28),
                             Inches(y + 0.3), Inches(0.52), Inches(0.4))
    tag.fill.solid(); tag.fill.fore_color.rgb = col
    tag.line.fill.background(); tag.shadow.inherit = False
    txt(s, x + 0.28, y + 0.37, 0.52, 0.3, num, size=14, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, x + 0.28, y + 0.92, cwid - 0.56, 0.4, title, size=17, bold=True)
    txt(s, x + 0.28, y + 1.42, cwid - 0.56, 1.9, body, size=12.5, color=TEXT2,
        line_spacing=1.3)
    card(s, x + 0.28, y + 3.15, cwid - 0.56, 0.68, fill=PRIMARY_SOFT, line=None)
    txt(s, x + 0.42, y + 3.28, cwid - 0.84, 0.5, "→  " + fix, size=11.5,
        bold=True, color=PRIMARY, line_spacing=1.2)
footer(s, 2)

# ================================================================ 3 streams
s = slide()
y = header(s, "The four streams",
           "Vertical ownership — each developer owns their slice from database to screen")
rows = [
    [("A", True, A_C), ("Platform & Security", True, TEXT),
     "Schema, migrations, auth, the scope guard, immutability core, CI/CD, then dashboard & reports",
     ("M0 · M1 · M2 · M6", True, TEXT2)],
    [("B", True, B_C), ("Masters & Clients", True, TEXT),
     "All 13 master screens, client master, Excel import, working calendar, workflow template designer",
     ("M3 → joins C", True, TEXT2)],
    [("C", True, C_C), ("Tickets & Ribbon", True, TEXT),
     "Ticket CRUD, detail page, cycles/reopen, comments, attachments, the Workflow Ribbon, Journey grid",
     ("M4", True, TEXT2)],
    [("D", True, D_C), ("Engines & Realtime", True, TEXT),
     "SLA & escalation scanners, mail engine, notification centre, WebSocket infrastructure, chat",
     ("M5 · M7", True, TEXT2)],
]
table(s, ML, y, CW, ["", "Stream", "Owns", "Milestones"], rows,
      [0.5, 2.1, 7.5, 1.83], row_h=0.72, head_h=0.4, fsize=12)
card(s, ML, y + 3.55, CW, 0.92, fill=PRIMARY_SOFT, line=None)
txt(s, ML + 0.3, y + 3.72, CW - 0.6, 0.6,
    "Why vertical, not 2 backend + 2 frontend:  a horizontal split makes every feature need two people "
    "to coordinate before it can ship — and idles the frontend developer whenever the backend one is behind.",
    size=13, color=PRIMARY, line_spacing=1.3)
footer(s, 3)

# ================================================================ 4 who
s = slide()
y = header(s, "Who should take which stream",
           "The assignment matters more than the plan — two of these four are load-bearing")
items = [
    (("A · Platform & Security", A_C),
     "Your strongest backend and infrastructure person.",
     "Everything else depends on their first six weeks, and the scope guard is the highest-risk component in the system. They are also the natural schema arbiter."),
    (("C · Tickets & Ribbon", C_C),
     "Your strongest all-rounder.",
     "The largest slice — roughly 40% of the product surface and the hardest UI in it. Stream B joins them from week 10."),
    (("B · Masters & Clients", B_C),
     "Fastest at CRUD and forms.",
     "Broad but shallow: many screens, few hard problems. The Excel import wizard is the one genuinely tricky piece."),
    (("D · Engines & Realtime", D_C),
     "Comfortable with async and scheduling.",
     "The most independent — workers, queues and sockets touch little of the request path. Least damaged by a slightly later start."),
]
cy = y
for (name, col), lead, body in items:
    card(s, ML, cy, CW, 1.03)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML), Inches(cy + 0.12),
                             Inches(0.06), Inches(0.79))
    bar.fill.solid(); bar.fill.fore_color.rgb = col
    bar.line.fill.background(); bar.shadow.inherit = False
    txt(s, ML + 0.32, cy + 0.17, 3.1, 0.32, name, size=15, bold=True, color=col)
    txt(s, ML + 0.32, cy + 0.55, 3.1, 0.32, lead, size=11.5, color=TEXT2)
    txt(s, ML + 3.6, cy + 0.22, CW - 3.95, 0.7, body, size=12.5, color=TEXT,
        line_spacing=1.28)
    cy += 1.16
footer(s, 4)

# ================================================================ 5 sprint 0
s = slide()
y = header(s, "Sprint 0 · weeks 1–2 — nobody waits",
           "Each stream has two weeks of work that depends on nothing anyone else is writing")
rows = [
    [("A", True, A_C),
     "Maven multi-module skeleton · docker-compose (MySQL, Redis, MinIO, Mailpit) · "
     "the complete Flyway baseline schema — all ~28 tables, triggers, indexes · CI pipeline · DB users and grants",
     ("Nothing", True, SUCCESS)],
    [("B", True, B_C),
     "All seed data — 6 roles + permission matrix, 11 task types, 4 priorities, statuses, "
     "3 workflow templates · JPA entities and repositories · the ticket fixture corpus",
     ("A's schema, day 3", False, TEXT2)],
    [("C", True, C_C),
     "React scaffold · design tokens from blueprint §12.1 · the shared component library — "
     "button, table, chip, modal, slide-over, toast, skeleton · Storybook · app shell",
     ("Nothing", True, SUCCESS)],
    [("D", True, D_C),
     "The OpenAPI contract for every endpoint in blueprint §13 · codegen pipeline · "
     "MSW mock server with realistic fixtures · the CI staleness check",
     ("Nothing", True, SUCCESS)],
]
table(s, ML, y, CW, ["", "Sprint 0 deliverable", "Depends on"], rows,
      [0.5, 9.7, 1.73], row_h=0.82, head_h=0.4, fsize=11.5)
card(s, ML, y + 3.85, CW, 0.82, fill=RGBColor(0xEC, 0xFD, 0xF5), line=None)
txt(s, ML + 0.3, y + 4.0, CW - 0.6, 0.55,
    "Why one person writes the schema:  28 interlocking tables authored by four people produce four naming "
    "conventions, four opinions on nullability, and a conflict on every migration file.",
    size=12.5, color=RGBColor(0x06, 0x5F, 0x46), line_spacing=1.3)
footer(s, 5)

# ================================================================ 6 decouplers
s = slide()
y = header(s, "Three decouplers keep the streams independent",
           "Without these, the dependency chain serialises the team — one person works, three wait")
cwid = (CW - 0.5) / 3
data = [
    ("Mock server", "owned by D · Sprint 0", B_C,
     "Every frontend feature is developed against MSW handlers generated from the OpenAPI spec.",
     "Stream C builds the entire ticket detail page — ribbon, tabs, journey grid — before a single ticket endpoint exists."),
    ("dev-noauth profile", "owned by A · day 10", A_C,
     "A Spring profile injecting a configurable fake principal: role, project list, reportee list.",
     "B, C and D develop against realistic scope behaviour four weeks before the real guard lands. Rejected outside local; disabled in CI."),
    ("Seed fixtures", "owned by B · Sprint 0", C_C,
     "200 tickets across 3 projects — varied stages, iterations, cycles and breach states.",
     "D tests the SLA scanner and C tests the ribbon before either feature exists. Build once, use everywhere."),
]
for i, (title, who, col, what, why) in enumerate(data):
    x = ML + i * (cwid + 0.25)
    card(s, x, y, cwid, 4.2)
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                             Inches(cwid), Inches(0.075))
    top.fill.solid(); top.fill.fore_color.rgb = col
    top.line.fill.background(); top.shadow.inherit = False
    txt(s, x + 0.28, y + 0.35, cwid - 0.56, 0.35, title, size=17, bold=True)
    txt(s, x + 0.28, y + 0.78, cwid - 0.56, 0.28, who.upper(), size=10,
        bold=True, color=col)
    txt(s, x + 0.28, y + 1.22, cwid - 0.56, 1.1, what, size=12.5, color=TEXT,
        line_spacing=1.3)
    txt(s, x + 0.28, y + 2.5, cwid - 0.56, 1.5, why, size=12, color=TEXT2,
        line_spacing=1.3)
footer(s, 6)

# ================================================================ 7 timeline
s = slide()
y = header(s, "Phase timeline · 18 weeks",
           "Consistent with the blueprint's own five-month estimate")
rows = [
    [("1–2", True, TEXT), "M0 scaffold, schema, CI", "Seed data, entities, fixtures",
     "Design system, Storybook", "OpenAPI contract, mocks"],
    [("3–5", True, TEXT), "M1 auth, JWT, rotation", "S-07–S-12 core masters",
     "S-19 create, S-17 list", "Outbox, scheduler, STOMP"],
    [("6–7", True, TEXT), "M1 scope guard + matrix", "S-32–S-34 clients + import",
     "S-20 detail, S-21 quick update", "M5 SLA + stage-SLA scanners"],
    [("8–9", True, TEXT), "M2 immutability, hash chain", "S-13/S-30 workflow designer",
     "M4 cycles, comments, attachments", "M5 escalation, mail engine"],
    [("10–11", True, TEXT), "M6 aggregation + dashboard", "joins C — ribbon UI",
     "M4 handoff, Journey grid", "M5 notifications, S-26"],
    [("12–14", True, TEXT), "M6 widgets, drill-down", "M6 reports + exports",
     "M4 S-22/23/24/29/31", "M7 chat, Ask Status"],
    [("15–16", True, TEXT), "Performance, S-16 audit", "S-28 Resource 360",
     "E2E walkthroughs A/B/C", "M7 finish, live ribbon"],
    [("17–18", True, TEXT), "Security review, go-live", "UAT fixes", "UAT fixes", "UAT fixes"],
]
table(s, ML, y, CW, ["Wk", "A — Platform", "B — Masters", "C — Tickets", "D — Engines"],
      rows, [0.72, 2.79, 2.79, 2.79, 2.84], row_h=0.4, head_h=0.4, fsize=10.5,
      hsize=11)
card(s, ML, y + 3.85, CW, 0.75, fill=RGBColor(0xFE, 0xF3, 0xC7), line=None)
txt(s, ML + 0.3, y + 4.0, CW - 0.6, 0.5,
    "The one hard dependency:  A's scope guard lands week 7. Until then B, C and D must NOT write their own "
    "filtering as a workaround — that becomes a permanent security hole.",
    size=12.5, color=RGBColor(0x92, 0x40, 0x0E), line_spacing=1.3)
footer(s, 7)

# ================================================================ 8 ownership
s = slide()
y = header(s, "Code ownership map",
           "Feature packaging, not layer packaging — so a stream's work lives where nobody else touches")
code_block(s, ML, y, 6.55, 4.35,
           "backend/\n"
           "  common/                    → A   shared DTOs, hashing, OpenAPI\n"
           "  domain/db/migration/       → A   Flyway — schema arbiter\n"
           "  api/security/              → A   filter chain, scope resolver\n"
           "  api/feature/auth|dashboard|reports/  → A\n"
           "  api/feature/masters|clients|imports/ → B\n"
           "  api/feature/workflow/      → B   (C consumes, B owns)\n"
           "  api/feature/tickets/       → C   CRUD, cycles, comments\n"
           "  api/feature/transitions/   → C   handoff, ribbon, journey\n"
           "  api/feature/notifications|chat/      → D\n"
           "  api/realtime/              → D   STOMP, interceptor, topics\n"
           "  worker/                    → D   (A owns hash verifier)\n"
           "\n"
           "frontend/src/\n"
           "  components/ui/             → C   additive changes only\n"
           "  components/ribbon/         → C\n"
           "  styles/tokens.css          → C   frozen after Sprint 0\n"
           "  api/generated/             → nobody — generated, never edited\n"
           "  features/auth|dashboard|reports/     → A\n"
           "  features/masters|clients/  → B\n"
           "  features/tickets/          → C\n"
           "  features/chat|notifications/         → D",
           size=10.5)
x2 = ML + 6.85
card(s, x2, y, CW - 6.85, 2.05, fill=PRIMARY_SOFT, line=None)
txt(s, x2 + 0.3, y + 0.24, CW - 7.45, 1.6,
    "The default Spring layout puts every controller in one folder.\n\n"
    "Four developers then edit that folder daily — which is exactly how "
    "a team of four ends up moving like a team of one.",
    size=13, color=PRIMARY, line_spacing=1.35)
card(s, x2, y + 2.25, CW - 6.85, 2.1)
txt(s, x2 + 0.3, y + 2.5, CW - 7.45, 0.3, "THE RULE", size=10.5, bold=True,
    color=TEXT2)
txt(s, x2 + 0.3, y + 2.88, CW - 7.45, 1.4,
    "Work only in your stream's paths.\n\n"
    "Touching another stream's directory needs that owner's sign-off — say so, "
    "rather than editing it quietly.",
    size=13, color=TEXT, line_spacing=1.35)
txt(s, x2, y + 4.55, CW - 6.85, 0.3, "This map becomes CODEOWNERS.", size=11.5,
    bold=True, color=TEXT2)
footer(s, 8)

# ================================================================ 9 hotspots
s = slide()
y = header(s, "Conflict hotspots — and the rule for each",
           "Seven files will otherwise conflict on most merges")
rows = [
    [("Flyway migrations", True, DANGER),
     "Sequential V14__ numbering collides whenever two developers add a migration the same day",
     ("Timestamp versioning: V20260812_1430__desc.sql  ·  never edit an applied file", True, TEXT)],
    [("Generated API client", True, WARNING),
     "A large generated tree that conflicts constantly and is meaningless to merge by hand",
     ("On any conflict take neither side — regenerate. .gitattributes merge=ours", False, TEXT)],
    [("Seed data", True, WARNING),
     "One shared seed file guarantees four-way conflicts",
     ("One seed file per stream, loaded by a manifest A owns", False, TEXT)],
    [("Security config", True, WARNING),
     "Every stream wants to register its own routes and permissions",
     ("A owns the filter chain; streams contribute per-feature @Configuration", False, TEXT)],
    [("Design tokens", True, INFO),
     "A stream adding its own colour drifts the palette across four features",
     ("Frozen after Sprint 0 — request a token from C, never add one", False, TEXT)],
    [("pom.xml / package.json", True, INFO),
     "Two developers adding different JSON libraries in the same week",
     ("Announce dependency changes before committing; regenerate lockfiles", False, TEXT)],
    [("Shared components", True, INFO),
     "C owns components/ui but all four consume it",
     ("Additive only. Storybook is the contract — not in Storybook, not shared", False, TEXT)],
]
table(s, ML, y, CW, ["File", "Why it conflicts", "The rule"], rows,
      [2.1, 4.5, 5.33], row_h=0.5, head_h=0.4, fsize=10.5)
footer(s, 9)

# ================================================================ 10 branch model
s = slide()
y = header(s, "Git branch model",
           "Developers push to branches · Claude merges · then it reaches main")
code_block(s, ML, y, 7.4, 2.6,
           "main                    protected · release-only · tagged\n"
           "                        Claude promotes — nobody pushes\n"
           "  └── develop           protected · integration\n"
           "                        Claude merges — nobody pushes\n"
           "       ├── feat/platform/jwt-refresh-rotation\n"
           "       ├── feat/masters/client-excel-import\n"
           "       ├── feat/tickets/ribbon-handoff-dialog\n"
           "       └── feat/engines/sla-scanner", size=12)
x2 = ML + 7.7
card(s, x2, y, CW - 7.7, 2.6, fill=PRIMARY_SOFT, line=None)
txt(s, x2 + 0.3, y + 0.26, CW - 8.3, 0.3, "WHY TWO PROTECTED BRANCHES",
    size=10.5, bold=True, color=PRIMARY)
txt(s, x2 + 0.3, y + 0.66, CW - 8.3, 1.8,
    "develop is where four streams actually integrate, and where conflicts surface.\n\n"
    "main only ever receives a develop that is already proven green.\n\n"
    "Merging four streams straight into main means resolving every conflict "
    "against your release branch.",
    size=12, color=PRIMARY, line_spacing=1.3)

txt(s, ML, y + 2.95, CW, 0.3, "NAMING  ·  <type>/<stream>/<slug>", size=11.5,
    bold=True, color=TEXT2)
card(s, ML, y + 3.35, 6.2, 1.35)
txt(s, ML + 0.3, y + 3.55, 5.6, 1.0,
    "type    feat · fix · chore · refactor · docs\n"
    "stream  platform · masters · tickets · engines\n"
    "slug    kebab-case, describes the change",
    size=12, font=MONO, color=TEXT, line_spacing=1.4)
card(s, ML + 6.5, y + 3.35, CW - 6.5, 1.35)
txt(s, ML + 6.8, y + 3.5, CW - 7.1, 0.3, "CONVENTIONAL COMMITS", size=10.5,
    bold=True, color=TEXT2)
txt(s, ML + 6.8, y + 3.85, CW - 7.1, 0.75,
    "feat(tickets): add handoff dialog with effort confirmation\n"
    "fix(engines): honour working calendar in stage-SLA check",
    size=11, font=MONO, color=TEXT, line_spacing=1.4)
footer(s, 10)

# ================================================================ 11 dev rules
s = slide()
y = header(s, "Rules for developers", "Seven rules that make the integration job tractable")
left = [
    ("Branch from develop", ", never from main."),
    ("Rebase on develop daily", " — a branch unrebased for a week is a conflict waiting to become someone else's problem."),
    ("Small pull requests", " — target under 400 changed lines. A 3,000-line PR cannot be meaningfully reviewed."),
    ("No merge commits inside a branch", " — rebase, so conflicts are resolved once rather than replayed."),
]
right = [
    ("Push at least once a day", ", even work in progress. Unpushed work is invisible work."),
    ("CI must be green", " before requesting integration. Claude does not merge a red branch."),
    ("Never commit", " .env, credentials, application-local.yml, target/, node_modules/, .DS_Store."),
]
bullets(s, ML, y + 0.1, 5.9, left, size=13, gap=1.0)
bullets(s, ML + 6.4, y + 0.1, 5.9, right, size=13, gap=1.0)
card(s, ML, y + 3.35, CW, 1.35, fill=RGBColor(0x1E, 0x1B, 0x33), line=None)
txt(s, ML + 0.35, y + 3.55, CW - 0.7, 0.3, "BRANCH PROTECTION  ·  main and develop",
    size=10.5, bold=True, color=RGBColor(0xA5, 0xB4, 0xFC))
txt(s, ML + 0.35, y + 3.95, CW - 0.7, 0.6,
    "Direct pushes blocked   ·   1 approval required   ·   status checks must pass   ·   "
    "linear history   ·   no force push   ·   no deletion",
    size=13, color=WHITE, line_spacing=1.3)
footer(s, 11)

# ================================================================ 12 integration
s = slide()
y = header(s, "The integration procedure", "What Claude runs — written down so it is repeatable, not ad hoc")
txt(s, ML, y, 6.4, 0.3, "PER BRANCH  ·  DAILY", size=10.5, bold=True, color=TEXT2)
code_block(s, ML, y + 0.36, 6.4, 2.5,
           "git fetch --all --prune\n"
           "git checkout develop && git pull\n"
           "\n"
           "git checkout feat/tickets/ribbon-handoff\n"
           "git rebase develop          # conflicts resolved here\n"
           "mvn -q verify\n"
           "npm --prefix frontend run build && npm test\n"
           "\n"
           "git checkout develop\n"
           "git merge --no-ff feat/tickets/ribbon-handoff\n"
           "git push origin develop", size=10.5)
txt(s, ML, y + 3.05, 6.4, 0.55,
    "--no-ff is deliberate: it keeps each feature visible as a unit, which matters "
    "when a milestone has to be bisected or reverted.",
    size=11.5, color=TEXT2, line_spacing=1.3)

x2 = ML + 6.75
txt(s, x2, y, CW - 6.75, 0.3, "PRE-MERGE CHECKLIST", size=10.5, bold=True, color=TEXT2)
card(s, x2, y + 0.36, CW - 6.75, 3.24)
checks = [
    "CI green on the branch",
    "No edits to another stream's paths without sign-off",
    "No changes to already-applied migrations",
    "Append-only tables still have no update or delete path",
    "New endpoints in the OpenAPI spec; client regenerated",
    "Tests exist; new routes have permission-matrix entries",
    "No secrets, no PII in logs, no stray debug output",
]
cy = y + 0.62
for c in checks:
    tick = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x2 + 0.28), Inches(cy + 0.05),
                              Inches(0.14), Inches(0.14))
    tick.fill.solid(); tick.fill.fore_color.rgb = SUCCESS
    tick.line.fill.background(); tick.shadow.inherit = False
    txt(s, x2 + 0.58, cy - 0.02, CW - 7.6, 0.4, c, size=12, color=TEXT,
        line_spacing=1.2)
    cy += 0.42
card(s, x2, y + 3.75, CW - 6.75, 0.95, fill=RGBColor(0xFE, 0xF3, 0xC7), line=None)
txt(s, x2 + 0.28, y + 3.93, CW - 7.3, 0.65,
    "On genuinely ambiguous conflicts — two streams changing the same business rule in "
    "different directions — the branch is handed back with the question, not guessed at.",
    size=11.5, color=RGBColor(0x92, 0x40, 0x0E), line_spacing=1.28)
footer(s, 12)

# ================================================================ 13 cadence
s = slide()
y = header(s, "Cadence and promotion to main", "")
rows = [
    [("Continuously", True, TEXT), "Developers push to their feature branches", ("Developers", False, TEXT2)],
    [("Daily · end of day", True, TEXT), "Claude integrates all ready branches into develop", ("Claude", False, PRIMARY)],
    [("Daily · next morning", True, TEXT), "Anything unmergeable is reported back with the specific conflict", ("Claude", False, PRIMARY)],
    [("Weekly / milestone end", True, TEXT), "Claude promotes develop → main and tags the release", ("Claude", False, PRIMARY)],
]
table(s, ML, y, CW, ["When", "What happens", "Who"], rows,
      [2.8, 7.6, 1.53], row_h=0.52, head_h=0.4, fsize=12.5)
txt(s, ML, y + 2.62, CW, 0.3, "PROMOTION", size=10.5, bold=True, color=TEXT2)
code_block(s, ML, y + 2.98, 7.0, 1.7,
           "git checkout develop && git pull\n"
           "mvn -q verify && npm --prefix frontend run build\n"
           "# smoke test against a fresh docker compose stack\n"
           "git checkout main && git merge --no-ff develop\n"
           "git tag -a v0.3.0 -m \"M3 — master data module\"\n"
           "git push origin main --tags", size=11)
card(s, ML + 7.3, y + 2.98, CW - 7.3, 1.7, fill=RGBColor(0xEC, 0xFD, 0xF5), line=None)
txt(s, ML + 7.6, y + 3.2, CW - 7.9, 1.3,
    "main should always be deployable.\n\n"
    "If it isn't, the promotion gate failed — and that is the thing to fix.",
    size=13.5, color=RGBColor(0x06, 0x5F, 0x46), line_spacing=1.35)
footer(s, 13)

# ================================================================ 14 bootstrap
s = slide()
y = header(s, "Repository bootstrap checklist",
           "The project is not yet a git repository, and gh is not installed on this machine")
left_items = [
    "git init, with develop as the working branch",
    ".gitignore — Java, Node, IDE, OS, secrets",
    ".gitattributes — generated client merge=ours, eol=lf",
    "Move blueprint, PLAN.md, TEAM-PLAN.md into docs/",
    "CODEOWNERS from the ownership map",
]
right_items = [
    "CONTRIBUTING.md — the rules, condensed to one page",
    "Pull request template",
    "Initial commit on main; branch develop; push both",
    "Create the remote, apply branch protection, add the team",
    ".github/workflows/ci.yml — build, test, staleness check",
]
cy = y
for i, it in enumerate(left_items):
    n = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(ML), Inches(cy), Inches(0.36), Inches(0.36))
    n.fill.solid(); n.fill.fore_color.rgb = PRIMARY
    n.line.fill.background(); n.shadow.inherit = False
    txt(s, ML, cy + 0.06, 0.36, 0.26, str(i + 1), size=12, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, ML + 0.52, cy + 0.04, 5.3, 0.5, it, size=13, color=TEXT, line_spacing=1.2)
    cy += 0.62
cy = y
for i, it in enumerate(right_items):
    x = ML + 6.3
    col = PRIMARY if i < 3 else TEXT2
    n = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(cy), Inches(0.36), Inches(0.36))
    n.fill.solid(); n.fill.fore_color.rgb = col
    n.line.fill.background(); n.shadow.inherit = False
    txt(s, x, cy + 0.06, 0.36, 0.26, str(i + 6), size=12, bold=True,
        color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, x + 0.52, cy + 0.04, 5.3, 0.5, it, size=13, color=TEXT, line_spacing=1.2)
    cy += 0.62
card(s, ML, y + 3.35, CW, 0.85, fill=PRIMARY_SOFT, line=None)
txt(s, ML + 0.3, y + 3.53, CW - 0.6, 0.6,
    "Steps 1–8 can be done locally today.   Step 9 needs the remote decided and either gh installed "
    "or the web UI.   Step 10 lands with Stream A's Sprint 0 work.",
    size=12.5, color=PRIMARY, line_spacing=1.3)
footer(s, 14)

# ================================================================ 15 DoD
s = slide()
y = header(s, "Definition of done", "A branch is ready for integration when all ten hold")
items = [
    "Feature works against the real backend, not only against mocks",
    "Unit tests for business logic; integration tests for any new endpoint",
    "New routes have permission-matrix entries for all six roles",
    "Migrations use timestamp versioning; no applied file was edited",
    "OpenAPI spec updated; client regenerated",
    "Storybook entry for any new shared component",
    "No new lint or compiler warnings",
    "Screens match the design tokens; keyboard navigable; ARIA labels present",
    "Rebased on current develop, CI green",
    "Touches only the stream's owned paths, or has the other owner's sign-off",
]
cy = y
for i, it in enumerate(items):
    x = ML if i < 5 else ML + 6.3
    yy = cy + (i % 5) * 0.62
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(yy),
                             Inches(0.32), Inches(0.32))
    box.fill.solid(); box.fill.fore_color.rgb = SURFACE
    box.line.color.rgb = SUCCESS; box.line.width = Pt(1.5)
    box.shadow.inherit = False
    try:
        box.adjustments[0] = 0.25
    except Exception:
        pass
    txt(s, x + 0.5, yy - 0.01, 5.4, 0.55, it, size=12.5, color=TEXT,
        line_spacing=1.2)
card(s, ML, y + 3.3, CW, 0.9, fill=RGBColor(0x1E, 0x1B, 0x33), line=None)
txt(s, ML + 0.35, y + 3.5, CW - 0.7, 0.55,
    "Developers never merge their own branch.   Push it, say it is ready, and it enters the "
    "integration queue.",
    size=14, color=WHITE, line_spacing=1.3)
footer(s, 15)

# ================================================================ 16 next
s = slide(SURFACE)
band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SW), Inches(SH))
band.fill.gradient()
band.fill.gradient_angle = 45.0
st = band.fill.gradient_stops
st[0].color.rgb = RGBColor(0x4F, 0x46, 0xE5); st[0].position = 0.0
st[1].color.rgb = RGBColor(0x1E, 0x1B, 0x4B); st[1].position = 1.0
band.line.fill.background(); band.shadow.inherit = False

txt(s, ML + 0.3, 0.85, CW, 0.6, "What to do right now", size=32, bold=True, color=WHITE)
line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML + 0.3), Inches(1.62),
                          Inches(1.3), Inches(0.04))
line.fill.solid(); line.fill.fore_color.rgb = RGBColor(0xA5, 0xB4, 0xFC)
line.line.fill.background(); line.shadow.inherit = False

steps = [
    ("1", "Assign the four streams", "Stream A must be your strongest backend person — they gate the other three for six weeks."),
    ("2", "Answer the directory question", "Does an employee directory or SSO already exist? If yes, users needs an external ID now, not after M3."),
    ("3", "Bootstrap the repository", "git init, .gitignore, docs/, CODEOWNERS, CONTRIBUTING.md, initial commit, branch develop. ~10 minutes."),
    ("4", "Create the remote and add the team", "Branch protection on both main and develop. Needs gh installed or the web UI."),
    ("5", "Kickoff — 90 minutes, whole team", "Walk blueprint §4A and §4 together. Four people need one shared mental model."),
    ("6", "Everyone runs their stream skill", "Sprint 0 begins. C and D start day 1; B starts day 3."),
]
cy = 2.05
for num, title, body in steps:
    n = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(ML + 0.3), Inches(cy),
                           Inches(0.42), Inches(0.42))
    n.fill.solid(); n.fill.fore_color.rgb = RGBColor(0xA5, 0xB4, 0xFC)
    n.line.fill.background(); n.shadow.inherit = False
    txt(s, ML + 0.3, cy + 0.09, 0.42, 0.28, num, size=13, bold=True,
        color=RGBColor(0x1E, 0x1B, 0x4B), align=PP_ALIGN.CENTER)
    txt(s, ML + 0.95, cy + 0.0, 3.5, 0.3, title, size=15, bold=True, color=WHITE)
    txt(s, ML + 4.6, cy + 0.03, CW - 4.9, 0.6, body, size=12.5,
        color=RGBColor(0xC7, 0xD2, 0xFE), line_spacing=1.25)
    cy += 0.74

txt(s, ML + 0.3, 6.75, CW, 0.3,
    "docs/TEAM-PLAN.md  ·  docs/GETTING-STARTED.md  ·  docs/streams/STREAM-A…D.md",
    size=11.5, color=RGBColor(0x81, 0x8C, 0xF8))

out = str(pathlib.Path(__file__).parents[2] / "docs/decks") + "/EduTrack-Team-Plan.pptx"
prs.save(out)
print("saved:", out)
print("slides:", len(prs.slides.__iter__.__self__._sldIdLst))
