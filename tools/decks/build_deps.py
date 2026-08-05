#!/usr/bin/env python3
"""EduTrack — Ownership & Dependency deck, generated from docs/DEPENDENCIES.md"""
import sys, pathlib
SP = str(pathlib.Path(__file__).parent)
sys.path.insert(0, SP)
from deck_lib import *
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

new_deck()
FT = "Ownership & Dependencies"
n = [0]
def pg(s):
    n[0] += 1; footer(s, n[0], FT)

OWN = {'A': ('Shivendra', A_C), 'B': ('Ayush', B_C), 'C': ('Divyansh', C_C), 'D': ('Debashis', D_C)}

def grad(s):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SW), Inches(SH))
    b.fill.gradient(); b.fill.gradient_angle = 45.0
    st = b.fill.gradient_stops
    st[0].color.rgb = RGBColor(0x4F,0x46,0xE5); st[0].position = 0.0
    st[1].color.rgb = RGBColor(0x1E,0x1B,0x4B); st[1].position = 1.0
    b.line.fill.background(); b.shadow.inherit = False

def bar(s, x, y, w, h, col, radius=True):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                            Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = col
    sh.line.fill.background(); sh.shadow.inherit = False
    if radius:
        try: sh.adjustments[0] = 0.22
        except Exception: pass
    return sh

def arrow_down(s, x, y1, y2, col, head=0.09):
    """vertical dependency arrow"""
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x-0.008), Inches(y1), Inches(0.016), Inches(y2-y1-head))
    ln.fill.solid(); ln.fill.fore_color.rgb = col; ln.line.fill.background(); ln.shadow.inherit = False
    tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x-0.055), Inches(y2-head), Inches(0.11), Inches(head))
    tri.rotation = 180
    tri.fill.solid(); tri.fill.fore_color.rgb = col; tri.line.fill.background(); tri.shadow.inherit = False

def arrow_right(s, x1, x2, y, col, head=0.11):
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x1), Inches(y-0.008), Inches(x2-x1-head), Inches(0.016))
    ln.fill.solid(); ln.fill.fore_color.rgb = col; ln.line.fill.background(); ln.shadow.inherit = False
    tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(x2-head), Inches(y-0.055), Inches(head), Inches(0.11))
    tri.rotation = 90
    tri.fill.solid(); tri.fill.fore_color.rgb = col; tri.line.fill.background(); tri.shadow.inherit = False

# ═══════════════════════════════════════════════════ 1 title
s = slide(SURFACE); grad(s)
txt(s, ML+0.3, 1.95, CW, 0.4, "EDUTRACK", size=15, bold=True, color=RGBColor(0xA5,0xB4,0xFC))
txt(s, ML+0.3, 2.45, CW, 1.5, "Ownership &\nDependency Map", size=42, bold=True, color=WHITE, line_spacing=1.08)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML+0.3), Inches(4.3), Inches(1.5), Inches(0.045))
ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0xA5,0xB4,0xFC); ln.line.fill.background(); ln.shadow.inherit = False
txt(s, ML+0.3, 4.62, CW-1.0, 1.1,
    "Who does what  ·  what blocks what  ·  what runs at the same time\n"
    "4 developers · 17 cross-stream dependencies · 18 weeks",
    size=15.5, color=RGBColor(0xC7,0xD2,0xFE), line_spacing=1.35)
txt(s, ML+0.3, 6.5, CW, 0.3, "docs/DEPENDENCIES.md", size=11.5, color=RGBColor(0x81,0x8C,0xF8), font=MONO)

# ═══════════════════════════════════════════════════ 2 the four streams
s = slide()
y = header(s, "Four vertical streams, one owner each",
           "Each owns their slice from database to screen — not a backend/frontend split")
cw = (CW - 0.75) / 4
data = [
 ('A', 'Platform & Security', 'Shivendra', 'A-001 … A-075',
  'Schema · migrations · auth · scope guard · immutability core · CI · dashboard · reports', A_C),
 ('B', 'Masters & Clients', 'Ayush', 'B-001 … B-063',
  'Master screens · client master · Excel import · working calendar · workflow designer', B_C),
 ('C', 'Tickets & Ribbon', 'Divyansh', 'C-001 … C-064',
  'Ticket CRUD · cycles · comments · attachments · the Workflow Ribbon · Journey grid', C_C),
 ('D', 'Engines & Realtime', 'Debashis', 'D-001 … D-059',
  'OpenAPI contract · SLA scanners · mail engine · notifications · WebSocket · chat', D_C),
]
for i,(k,name,owner,ids,scope,col) in enumerate(data):
    x = ML + i*(cw+0.25)
    card(s, x, y, cw, 3.5)
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(cw), Inches(0.075))
    top.fill.solid(); top.fill.fore_color.rgb = col; top.line.fill.background(); top.shadow.inherit = False
    txt(s, x+0.28, y+0.32, 0.5, 0.4, k, size=25, bold=True, color=col)
    txt(s, x+0.28, y+0.85, cw-0.56, 0.5, name, size=14.5, bold=True, line_spacing=1.15)
    txt(s, x+0.28, y+1.45, cw-0.56, 0.28, owner, size=13, bold=True, color=col)
    txt(s, x+0.28, y+1.76, cw-0.56, 0.25, ids, size=10.5, color=TEXT2, font=MONO)
    txt(s, x+0.28, y+2.15, cw-0.56, 1.2, scope, size=11, color=TEXT2, line_spacing=1.3)
card(s, ML, y+3.75, CW, 0.75, fill=PRIMARY_SOFT, line=None)
txt(s, ML+0.35, y+3.9, CW-0.7, 0.5,
    "Only 17 dependencies cross a stream boundary. Everything else is internal to one developer — "
    "which is the entire point of splitting the work this way.",
    size=13, color=PRIMARY, line_spacing=1.3)
pg(s)

# ═══════════════════════════════════════════════════ 3 four dates
s = slide()
y = header(s, "Four dates decide whether this works",
           "Everything else can slip a few days. These cannot — they are what makes four people work in parallel")
rows = [
 [("Day 5", True, DANGER), ("D-004", True, D_C), "MSW mock server", ("Debashis", False, TEXT),
  ("Divyansh's entire frontend stalls", True, TEXT)],
 [("Day 10", True, DANGER), ("A-012", True, A_C), "dev-noauth profile", ("Shivendra", False, TEXT),
  ("All three others stop", True, TEXT)],
 [("Day 10", True, DANGER), ("B-007", True, B_C), "Ticket fixture corpus", ("Ayush", False, TEXT),
  ("No data to test ribbon or SLA against", True, TEXT)],
 [("Week 3", True, WARNING), ("B-024", True, B_C), "Working-hours service", ("Ayush", False, TEXT),
  ("Debashis cannot start the SLA engine at all", True, TEXT)],
]
table(s, ML, y, CW, ["Due", "Task", "What it is", "Owner", "If it is late"], rows,
      [1.1, 1.0, 3.1, 1.9, 4.83], row_h=0.62, head_h=0.42, fsize=12.5)
card(s, ML, y+3.2, CW, 1.3, fill=RGBColor(0x1E,0x1B,0x33), line=None)
txt(s, ML+0.35, y+3.4, CW-0.7, 0.3, "WHY THESE FOUR AND NOT OTHERS", size=10.5, bold=True, color=RGBColor(0xA5,0xB4,0xFC))
txt(s, ML+0.35, y+3.78, CW-0.7, 0.65,
    "Each one unblocks somebody else rather than advancing its own stream. They are the decouplers — "
    "mocks let the frontend run ahead of the API, dev-noauth lets three people work before the security spine lands, "
    "and fixtures let two people test against data that does not exist yet.",
    size=12.5, color=WHITE, line_spacing=1.32)
pg(s)

# ═══════════════════════════════════════════════════ 4 dependency register 1
s = slide()
y = header(s, "The 17 cross-stream dependencies · 1 of 2",
           "Every edge where one developer waits on another")
rows = [
 [("1", False, TEXT2), ("A-002", True, A_C), "docker-compose", ("all three", True, TEXT), ("Day 1", True, SUCCESS)],
 [("2", False, TEXT2), ("A-003…A-009", True, A_C), "Baseline schema", ("Ayush — seeds & entities", False, TEXT), ("Day 3", True, TEXT)],
 [("3", False, TEXT2), ("A-012", True, A_C), "dev-noauth profile", ("all three — any authenticated endpoint", True, TEXT), ("Day 10", True, DANGER)],
 [("4", False, TEXT2), ("A-034", True, A_C), "ScopeResolver", ("Debashis — socket authorisation mirrors it", False, TEXT), ("Week 7", False, TEXT2)],
 [("5", False, TEXT2), ("A-040…A-042", True, A_C), "Append-only + hash chain", ("Divyansh — transitions write to those tables", True, TEXT), ("Week 9", True, WARNING)],
 [("6", False, TEXT2), ("B-001…B-004", True, B_C), "Seed data", ("all three — realistic local data", False, TEXT), ("Day 7", False, TEXT2)],
 [("7", False, TEXT2), ("B-007", True, B_C), "Fixture corpus", ("Divyansh + Debashis — ribbon & SLA testing", True, TEXT), ("Day 10", True, DANGER)],
 [("8", False, TEXT2), ("B-021", True, B_C), "Priority master", ("Divyansh — priority dropdown", False, TEXT), ("Week 5", False, TEXT2)],
 [("9", False, TEXT2), ("B-022", True, B_C), "Notification templates", ("Debashis — mail engine", False, TEXT), ("Week 8", False, TEXT2)],
]
table(s, ML, y, CW, ["#", "Blocker", "What it is", "Who waits on it", "Needed"], rows,
      [0.5, 1.75, 2.75, 5.4, 1.53], row_h=0.42, head_h=0.42, fsize=11)
pg(s)

# ═══════════════════════════════════════════════════ 5 dependency register 2
s = slide()
y = header(s, "The 17 cross-stream dependencies · 2 of 2", "")
rows = [
 [("10", False, TEXT2), ("B-024", True, B_C), "Working-hours service", ("Debashis — every SLA calculation · Divyansh — stage durations", True, TEXT), ("Week 3", True, DANGER)],
 [("11", False, TEXT2), ("B-039…B-043", True, B_C), "Workflow templates", ("Divyansh — transitions need stage definitions", True, TEXT), ("Week 9", True, WARNING)],
 [("12", False, TEXT2), ("B-025…B-029", True, B_C), "Client master", ("Divyansh — client dropdown · Shivendra — client widget", False, TEXT), ("Week 7", False, TEXT2)],
 [("13", False, TEXT2), ("C-042…C-049", True, C_C), "Transition service", ("Shivendra — ribbon widgets · Debashis — live push", True, TEXT), ("Week 11", False, TEXT2)],
 [("14", False, TEXT2), ("C-003", True, C_C), "Component library", ("all three — every screen in every stream", True, TEXT), ("Day 7", False, TEXT2)],
 [("15", False, TEXT2), ("D-001", True, D_C), "OpenAPI contract", ("all three — all API work", True, TEXT), ("Week 1", True, TEXT)],
 [("16", False, TEXT2), ("D-004", True, D_C), "MSW mock server", ("Divyansh — the entire frontend before APIs exist", True, TEXT), ("Day 5", True, DANGER)],
 [("17", False, TEXT2), ("D-012…D-014", True, D_C), "STOMP infrastructure", ("Divyansh — live ribbon advance", False, TEXT), ("Week 10", False, TEXT2)],
]
table(s, ML, y, CW, ["#", "Blocker", "What it is", "Who waits on it", "Needed"], rows,
      [0.5, 1.75, 2.75, 5.4, 1.53], row_h=0.44, head_h=0.42, fsize=11)
card(s, ML, y+4.05, CW, 0.62, fill=PRIMARY_SOFT, line=None)
txt(s, ML+0.35, y+4.18, CW-0.7, 0.4,
    "Shivendra is upstream of everyone and blocked by nobody. Ayush carries four of the seventeen edges — "
    "more than anyone else.",
    size=12.5, color=PRIMARY)
pg(s)

# ═══════════════════════════════════════════════════ 6 critical path
s = slide()
y = header(s, "The critical path", "Eight links. Delay anywhere on this chain delays go-live one-for-one")
chain = [
 ("A-003…A-009", "Schema", "wk 1", "Shivendra", A_C),
 ("A-008", "Immutability triggers", "wk 2", "Shivendra", A_C),
 ("A-040…A-042", "Append-only + hash chain", "wk 8–9", "Shivendra", A_C),
 ("C-042…C-049", "Transition service", "wk 10–11", "Divyansh", C_C),
 ("C-055…C-058", "Journey roll-up", "wk 11", "Divyansh", C_C),
 ("A-058", "Ribbon dashboard widgets", "wk 12–14", "Shivendra", A_C),
 ("E2E", "Walkthrough A end to end", "wk 15–16", "Divyansh", C_C),
 ("Go-live", "Hardening, UAT", "wk 17–18", "all four", PRIMARY),
]
bw = (CW - 0.1) / 8
for i,(tid,name,wk,who,col) in enumerate(chain):
    x = ML + i*bw
    card(s, x+0.04, y, bw-0.16, 2.5)
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x+0.04), Inches(y), Inches(bw-0.16), Inches(0.06))
    top.fill.solid(); top.fill.fore_color.rgb = col; top.line.fill.background(); top.shadow.inherit = False
    txt(s, x+0.18, y+0.24, bw-0.44, 0.3, tid, size=10, bold=True, color=col, font=MONO)
    txt(s, x+0.18, y+0.58, bw-0.44, 0.9, name, size=12, bold=True, line_spacing=1.2)
    txt(s, x+0.18, y+1.72, bw-0.44, 0.25, wk, size=10.5, color=TEXT2)
    txt(s, x+0.18, y+1.98, bw-0.44, 0.25, who, size=10, color=col, bold=True)
    if i < 7:
        arrow_right(s, x+bw-0.14, x+bw+0.04, y+1.25, RGBColor(0x9C,0xA3,0xAF), head=0.09)
card(s, ML, y+2.75, 5.85, 1.6, fill=RGBColor(0xFE,0xF3,0xC7), line=None)
txt(s, ML+0.32, y+2.95, 5.25, 0.3, "THE MOST FRAGILE HANDOFF", size=10.5, bold=True, color=RGBColor(0x92,0x40,0x0E))
txt(s, ML+0.32, y+3.32, 5.25, 1.0,
    "A-042 → C-042, week 9. Divyansh cannot start the transition service until the hash chain is finished and proven, "
    "because his service writes into those append-only tables.",
    size=12, color=RGBColor(0x92,0x40,0x0E), line_spacing=1.3)
card(s, ML+6.15, y+2.75, CW-6.15, 1.6, fill=RGBColor(0xEC,0xFD,0xF5), line=None)
txt(s, ML+6.47, y+2.95, CW-6.75, 0.3, "OFF THE CRITICAL PATH — HAS SLACK", size=10.5, bold=True, color=RGBColor(0x06,0x5F,0x46))
txt(s, ML+6.47, y+3.32, CW-6.75, 1.0,
    "Client master · Excel import · reports · chat · notification centre · browser push.\n"
    "These can absorb delay without moving go-live.",
    size=12, color=RGBColor(0x06,0x5F,0x46), line_spacing=1.3)
pg(s)

# ═══════════════════════════════════════════════════ 7 swimlanes
s = slide()
y = header(s, "18 weeks, four lanes",
           "Bars are work. Arrows are the moments one developer hands something to another")
L = ML + 1.15          # lane left
LW = CW - 1.3          # lane width
def wx(w):             # week -> x  (weeks 1..18)
    return L + (w-1)/18.0 * LW
def ww(a,b):
    return wx(b+1) - wx(a)

# week ruler
for w in range(1,19):
    if w % 2 == 1:
        txt(s, wx(w), y, 0.5, 0.22, str(w), size=9, color=TEXT2)
txt(s, ML, y, 1.0, 0.22, "WEEK", size=9, bold=True, color=TEXT2)

LANE_H = 0.52
GAP = 0.30
lanes = [
 ('A', [ (1,2,'schema · triggers · CI',A_C), (3,7,'auth · scope guard · matrix',A_C),
         (8,9,'immutability core',A_C), (10,16,'dashboard · reports',A_C), (17,18,'hardening',PRIMARY) ]),
 ('B', [ (1,2,'seeds · entities · fixtures',B_C), (3,9,'masters · clients · import · workflow',B_C),
         (10,11,'ribbon UI with C',B_C), (12,16,'reports · Resource 360',B_C), (17,18,'UAT',PRIMARY) ]),
 ('C', [ (1,2,'design system',C_C), (3,9,'create · list · detail · comments',C_C),
         (10,11,'transitions · ribbon',C_C), (12,14,'journey · stage queue',C_C),
         (15,16,'E2E',C_C), (17,18,'UAT',PRIMARY) ]),
 ('D', [ (1,2,'contract · mocks',D_C), (3,5,'outbox · STOMP',D_C),
         (6,11,'SLA · escalation · mail',D_C), (12,16,'chat · live ribbon',D_C), (17,18,'UAT',PRIMARY) ]),
]
top0 = y + 0.32
for li,(k,blocks) in enumerate(lanes):
    ly = top0 + li*(LANE_H+GAP)
    owner, col = OWN[k]
    txt(s, ML, ly+0.11, 1.1, 0.3, f"{k} · {owner}", size=11, bold=True, color=col)
    band = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(L), Inches(ly), Inches(LW), Inches(LANE_H))
    band.fill.solid(); band.fill.fore_color.rgb = SUBTLE; band.line.fill.background(); band.shadow.inherit = False
    for (a,b,label,c) in blocks:
        bx, bw2 = wx(a), ww(a,b)
        bar(s, bx+0.02, ly+0.05, bw2-0.04, LANE_H-0.1, c)
        txt(s, bx+0.12, ly+0.16, bw2-0.24, 0.28, label, size=9, color=WHITE, bold=True)

# dependency arrows between lanes
def lane_y(i, bottom=True):
    ly = top0 + i*(LANE_H+GAP)
    return ly + (LANE_H if bottom else 0)
# D-004 mocks (D, wk1) -> C frontend (C, wk3): D lane idx 3 -> C idx 2  (upward: draw from C bottom to D? simpler: vertical between lanes)
arrow_down(s, wx(3), lane_y(0), lane_y(1, False), RGBColor(0x9C,0xA3,0xAF))   # A schema -> B entities
arrow_down(s, wx(4), lane_y(2), lane_y(3, False), RGBColor(0x9C,0xA3,0xAF))   # C -> D  (mocks direction shown in note)
arrow_down(s, wx(10), lane_y(0), lane_y(2, False), DANGER)                     # A-042 -> C-042
arrow_down(s, wx(10.6), lane_y(1), lane_y(2, False), WARNING)                  # B-043 -> C-042
arrow_down(s, wx(12), lane_y(2), lane_y(3, False), RGBColor(0x9C,0xA3,0xAF))   # C transitions -> D live push

txt(s, ML, top0 + 4*(LANE_H+GAP) + 0.12, CW, 0.3,
    "Red and amber arrows at week 10 are the convergence: the transition service needs Shivendra's hash chain AND "
    "Ayush's stage definitions before it can start.", size=11.5, color=TEXT2)
card(s, ML, SH-1.35, CW, 0.62, fill=PRIMARY_SOFT, line=None)
txt(s, ML+0.35, SH-1.22, CW-0.7, 0.4,
    "Weeks 3–7 and 12–16 are fully parallel — nobody blocks anybody. Weeks 1–2 and 8–11 carry every handoff.",
    size=12.5, color=PRIMARY)
pg(s)

# ═══════════════════════════════════════════════════ 8 convergence
s = slide()
y = header(s, "Week 9 — the one convergence to protect",
           "Two upstream owners, one downstream task, sitting on the critical path")
# two sources
card(s, ML+1.1, y+0.15, 4.0, 1.35)
tb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML+1.1), Inches(y+0.15), Inches(4.0), Inches(0.07))
tb.fill.solid(); tb.fill.fore_color.rgb = A_C; tb.line.fill.background(); tb.shadow.inherit = False
txt(s, ML+1.4, y+0.42, 3.5, 0.3, "A-040 … A-042", size=13, bold=True, color=A_C, font=MONO)
txt(s, ML+1.4, y+0.75, 3.5, 0.55, "Append-only services + per-ticket hash chain\nShivendra", size=11.5, color=TEXT2, line_spacing=1.3)

card(s, ML+6.9, y+0.15, 4.0, 1.35)
tb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML+6.9), Inches(y+0.15), Inches(4.0), Inches(0.07))
tb.fill.solid(); tb.fill.fore_color.rgb = B_C; tb.line.fill.background(); tb.shadow.inherit = False
txt(s, ML+7.2, y+0.42, 3.5, 0.3, "B-039 … B-043", size=13, bold=True, color=B_C, font=MONO)
txt(s, ML+7.2, y+0.75, 3.5, 0.55, "Workflow templates + stage definitions\nAyush", size=11.5, color=TEXT2, line_spacing=1.3)

arrow_down(s, ML+3.1, y+1.55, y+2.35, A_C)
arrow_down(s, ML+8.9, y+1.55, y+2.35, B_C)

card(s, ML+3.0, y+2.4, 6.0, 1.35, fill=PRIMARY_SOFT, line=None)
txt(s, ML+3.35, y+2.65, 5.4, 0.3, "C-042 … C-049", size=13, bold=True, color=PRIMARY, font=MONO)
txt(s, ML+3.35, y+2.98, 5.4, 0.6, "Transition service — the Workflow Ribbon\nDivyansh · week 10–11 · ON THE CRITICAL PATH",
    size=11.5, color=PRIMARY, line_spacing=1.3)

card(s, ML, y+4.0, CW, 1.15, fill=RGBColor(0xFE,0xF3,0xC7), line=None)
txt(s, ML+0.35, y+4.18, CW-0.7, 0.3, "IF EITHER IS LATE", size=10.5, bold=True, color=RGBColor(0x92,0x40,0x0E))
txt(s, ML+0.35, y+4.52, CW-0.7, 0.5,
    "Give Divyansh C-059…C-064 — History and Attachments tabs, stage queue, bulk reassignment, ticket links. "
    "They depend on nothing and keep him productive while the blockers clear.",
    size=12.5, color=RGBColor(0x92,0x40,0x0E), line_spacing=1.3)
pg(s)

# ═══════════════════════════════════════════════════ 9 blocked
s = slide()
y = header(s, "If you are blocked, do this", "Nobody should ever idle — every stream has independent fallback work")
rows = [
 [("Shivendra", True, A_C), ("— nothing blocks him", False, TEXT2), ("He is upstream of everyone", False, TEXT2)],
 [("Ayush", True, B_C), "A's schema, days 1–2", ("Design the seed content — permission matrix, 11 task types, 4 priorities, 3 workflow templates. It comes from the blueprint, not the schema", False, TEXT)],
 [("Ayush", True, B_C), "A's schema, later", ("B-030…B-038 Excel import engine — pure POI work, touches no EduTrack table", False, TEXT)],
 [("Divyansh", True, C_C), "D's mocks, days 1–4", ("C-002…C-004 design tokens, component library, Storybook — no data needed", False, TEXT)],
 [("Divyansh", True, C_C), "A-042 or B-043, week 9", ("C-059…C-064 tabs, stage queue, bulk reassign, ticket links", False, TEXT)],
 [("Debashis", True, D_C), "B-024 working-hours, week 3", ("D-010…D-015 outbox, ShedLock, STOMP, channel interceptor", False, TEXT)],
 [("Debashis", True, D_C), "C's transitions, week 11", ("D-040…D-046 notification centre, preference matrix, offline queue", False, TEXT)],
]
table(s, ML, y, CW, ["Owner", "Blocked on", "Pick up instead"], rows,
      [1.7, 2.9, 7.33], row_h=0.52, head_h=0.42, fsize=11.5)
pg(s)

# ═══════════════════════════════════════════════════ 10 false deps
s = slide()
y = header(s, "Three habits that prevent imaginary blockers",
           "Each of these has stalled real teams on this exact shape of project")
items = [
 ("Never wait for an API — use the mocks", D_C,
  "D-004 exists so the full ticket detail page, ribbon included, can be built before a single ticket endpoint is written. "
  "When the real endpoint lands, the only change is a flag."),
 ("Never wait for auth — use dev-noauth", A_C,
  "And never write your own filtering as a workaround. That is how a temporary shortcut becomes a permanent security hole — "
  "the top risk in blueprint §17."),
 ("Never write your own date maths", B_C,
  "Every SLA, duration, breach and utilisation figure routes through B-024. Four private implementations of \"working hours\" "
  "produce four different answers, and the disagreement surfaces in a client dispute."),
]
cy = y
for t,col,b in items:
    card(s, ML, cy, CW, 1.32)
    bar2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML), Inches(cy+0.14), Inches(0.06), Inches(1.04))
    bar2.fill.solid(); bar2.fill.fore_color.rgb = col; bar2.line.fill.background(); bar2.shadow.inherit = False
    txt(s, ML+0.34, cy+0.26, CW-0.7, 0.35, t, size=16, bold=True, color=col)
    txt(s, ML+0.34, cy+0.72, CW-0.7, 0.5, b, size=12.5, color=TEXT2, line_spacing=1.3)
    cy += 1.46
pg(s)

# ═══════════════════════════════════════════════════ 11 closing
s = slide(SURFACE); grad(s)
txt(s, ML+0.3, 0.95, CW, 0.55, "How to use this", size=32, bold=True, color=WHITE)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML+0.3), Inches(1.72), Inches(1.3), Inches(0.04))
ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0xA5,0xB4,0xFC); ln.line.fill.background(); ln.shadow.inherit = False
pts = [
 ("Daily", "Check whether anything you own appears in the 17-edge register. If it does, somebody is waiting on you."),
 ("When blocked", "Go to the fallback table. There is always independent work — idling is never the right answer."),
 ("At milestone boundaries", "Review the register when develop is promoted to main. Tasks move between streams; the map must move with them."),
 ("When adding a task", "If it crosses a stream boundary, add it to §2 in the same pull request."),
]
cy = 2.15
for k,v in pts:
    txt(s, ML+0.3, cy, 3.2, 0.3, k, size=15, bold=True, color=RGBColor(0xA5,0xB4,0xFC))
    txt(s, ML+3.7, cy+0.02, CW-4.0, 0.7, v, size=13, color=WHITE, line_spacing=1.3)
    cy += 0.92
card(s, ML+0.3, 6.0, CW-0.6, 0.72, fill=RGBColor(0xFF,0xFF,0xFF), line=None)
txt(s, ML+0.6, 6.16, CW-1.2, 0.45,
    "A dependency register that has drifted is worse than none — people trust it and are wrong.",
    size=13.5, bold=True, color=PRIMARY)
txt(s, ML+0.3, 6.95, CW, 0.3, "docs/DEPENDENCIES.md  ·  docs/streams/STREAM-A…D.md  ·  docs/TEAM-PLAN.md",
    size=11.5, color=RGBColor(0x81,0x8C,0xF8))

save(str(pathlib.Path(__file__).parents[2] / "docs/decks") + "/EduTrack-Dependencies.pptx")
