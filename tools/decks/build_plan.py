#!/usr/bin/env python3
"""EduTrack — Implementation Plan deck, generated from PLAN.md"""
import sys, pathlib
sys.path.insert(0, str(pathlib.Path(__file__).parent))
from deck_lib import *
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

new_deck()
FT = "Implementation Plan"
n = [0]
def pg(s):
    n[0] += 1
    footer(s, n[0], FT)

def grad(s):
    b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(SW), Inches(SH))
    b.fill.gradient(); b.fill.gradient_angle = 45.0
    st = b.fill.gradient_stops
    st[0].color.rgb = RGBColor(0x4F,0x46,0xE5); st[0].position = 0.0
    st[1].color.rgb = RGBColor(0x1E,0x1B,0x4B); st[1].position = 1.0
    b.line.fill.background(); b.shadow.inherit = False

# ---------------------------------------------------------------- 1 title
s = slide(SURFACE); grad(s)
txt(s, ML+0.3, 2.05, CW, 0.4, "EDUTRACK", size=15, bold=True, color=RGBColor(0xA5,0xB4,0xFC))
txt(s, ML+0.3, 2.55, CW, 1.6, "Implementation Plan", size=44, bold=True, color=WHITE, line_spacing=1.08)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML+0.3), Inches(4.0), Inches(1.5), Inches(0.045))
ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0xA5,0xB4,0xFC)
ln.line.fill.background(); ln.shadow.inherit = False
txt(s, ML+0.3, 4.35, CW-1.0, 1.2,
    "Java 21 · Spring Boot 3.3 · MySQL 8.4 · React 18 + TypeScript\n"
    "PostgreSQL → MySQL translation · milestones M0–M7 · 18 weeks",
    size=15.5, color=RGBColor(0xC7,0xD2,0xFE), line_spacing=1.35)
txt(s, ML+0.3, 6.45, CW, 0.3, "Build specification — companion to the Product Blueprint v1.2",
    size=11.5, color=RGBColor(0x81,0x8C,0xF8))

# ---------------------------------------------------------------- 2 how to read
s = slide()
y = header(s, "Two documents, two authorities",
           "Where they disagree, each wins in its own domain — and nothing is silent")
card(s, ML, y, 5.8, 2.2, fill=PRIMARY_SOFT, line=None)
txt(s, ML+0.32, y+0.26, 5.2, 0.32, "THE BLUEPRINT", size=11, bold=True, color=PRIMARY)
txt(s, ML+0.32, y+0.68, 5.2, 0.4, "The product specification", size=17, bold=True, color=PRIMARY)
txt(s, ML+0.32, y+1.2, 5.2, 0.85,
    "What the system does. Every screen, every field, every rule.\n"
    "Wins on BEHAVIOUR.",
    size=13, color=PRIMARY, line_spacing=1.3)
card(s, ML+6.15, y, 5.78, 2.2, fill=RGBColor(0xEC,0xFD,0xF5), line=None)
txt(s, ML+6.47, y+0.26, 5.2, 0.32, "THIS PLAN", size=11, bold=True, color=RGBColor(0x06,0x5F,0x46))
txt(s, ML+6.47, y+0.68, 5.2, 0.4, "The build specification", size=17, bold=True, color=RGBColor(0x06,0x5F,0x46))
txt(s, ML+6.47, y+1.2, 5.2, 0.85,
    "What stack, what order, what has to be decided first.\n"
    "Wins on IMPLEMENTATION.",
    size=13, color=RGBColor(0x06,0x5F,0x46), line_spacing=1.3)
card(s, ML, y+2.5, CW, 1.95, fill=SURFACE)
txt(s, ML+0.35, y+2.7, CW-0.7, 0.3, "WHY THIS MATTERS HERE", size=11, bold=True, color=TEXT2)
txt(s, ML+0.35, y+3.08, CW-0.7, 1.2,
    "The blueprint recommends NestJS + PostgreSQL. We are building on Java + MySQL.\n\n"
    "That means its DDL, its BullMQ / Socket.IO / SheetJS choices and its shared-Zod-schema guidance do not apply "
    "as written. Every substitution is listed explicitly — none are left for a developer to discover mid-sprint.",
    size=13.5, color=TEXT, line_spacing=1.35)
pg(s)

# ---------------------------------------------------------------- 3 spine
s = slide()
y = header(s, "Three things are the architectural spine",
           "Everything else hangs off them, and each is expensive to retrofit")
spine = [
    ("01", "The append-only history model", A_C,
     "ticket_history, ticket_effort_logs and ticket_stage_transitions are insert-only, hash-chained, enforced at four "
     "independent layers — service, DB grants, DB triggers, and absent HTTP routes.",
     "Built late, every write path already written has to be re-audited."),
    ("02", "The Workflow Ribbon", C_C,
     "Eight configurable stages, cycle_no and iteration_no as two orthogonal counters, per-stage-per-resource effort "
     "attribution, and the active-vs-idle split that makes it more than decoration.",
     "The product differentiator, and the hardest UI in the system."),
    ("03", "The row-scope guard", D_C,
     "Every ticket query rewritten server-side with a mandatory scope predicate. Developer sees assigned_to = me; "
     "PM sees their projects; out-of-scope IDs return 404, not 403.",
     "One leak here and the system fails its first audit."),
]
cwid = (CW - 0.5) / 3
for i, (num, t, col, body, why) in enumerate(spine):
    x = ML + i * (cwid + 0.25)
    card(s, x, y, cwid, 4.2)
    tag = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x+0.28), Inches(y+0.3), Inches(0.52), Inches(0.4))
    tag.fill.solid(); tag.fill.fore_color.rgb = col
    tag.line.fill.background(); tag.shadow.inherit = False
    txt(s, x+0.28, y+0.37, 0.52, 0.3, num, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, x+0.28, y+0.92, cwid-0.56, 0.62, t, size=16.5, bold=True, line_spacing=1.15)
    txt(s, x+0.28, y+1.68, cwid-0.56, 1.6, body, size=12, color=TEXT2, line_spacing=1.3)
    card(s, x+0.28, y+3.28, cwid-0.56, 0.72, fill=RGBColor(0xFE,0xF2,0xF2), line=None)
    txt(s, x+0.42, y+3.42, cwid-0.84, 0.55, why, size=11, bold=True,
        color=RGBColor(0x99,0x1B,0x1B), line_spacing=1.22)
pg(s)

# ---------------------------------------------------------------- 4 stack
s = slide()
y = header(s, "The stack", "")
rows = [
    [("Language · Framework", True, TEXT), "Java 21 LTS · Spring Boot 3.3+", ("Virtual threads for the IMAP poller and mail workers", False, TEXT2)],
    [("Persistence", True, TEXT), "Spring Data JPA + Hibernate 6 · Flyway", ("JdbcTemplate for reporting queries; migrations hand-written SQL", False, TEXT2)],
    [("Database", True, TEXT), "MySQL 8.4 · InnoDB · utf8mb4", ("All timestamps DATETIME(6) in UTC", False, TEXT2)],
    [("Cache · Scheduler", True, TEXT), "Redis 7 · @Scheduled + ShedLock", ("Token registry, rate limits; ShedLock stops double-firing", False, TEXT2)],
    [("Async · Realtime", True, TEXT), "Transactional outbox · Spring WebSocket + STOMP", ("Redis relay for multi-instance fan-out", False, TEXT2)],
    [("Storage · Excel · Mail", True, TEXT), "MinIO/S3 · Apache POI · Spring Mail + Thymeleaf", ("POI streaming, not the DOM reader", False, TEXT2)],
    [("Security", True, TEXT), "Argon2id · Nimbus JOSE JWT", ("64 MB memory, 3 iterations", False, TEXT2)],
    [("Frontend", True, TEXT), "React 18 + TS + Vite · Tailwind + shadcn/ui", ("TanStack Query · Zustand · RHF + Zod · Recharts", False, TEXT2)],
    [("Contract", True, TEXT), "springdoc-openapi → generated TS client", ("CI fails if the committed client is stale", False, TEXT2)],
    [("Testing", True, TEXT), "JUnit 5 · Testcontainers · MockMvc · ArchUnit", ("ArchUnit enforces the scope-guard rule structurally", False, TEXT2)],
]
table(s, ML, y, CW, ["Layer", "Choice", "Note"], rows,
      [2.7, 4.4, 4.83], row_h=0.34, head_h=0.4, fsize=11)
pg(s)

# ---------------------------------------------------------------- 5 substitutions
s = slide()
y = header(s, "Five Node choices with no Java equivalent",
           "Each substitution preserves the behaviour the blueprint specifies")
rows = [
    [("BullMQ", True, DANGER), ("Transactional outbox + @Scheduled + ShedLock", True, TEXT),
     "email_log already IS an outbox — status, retry_count, provider ID. Claim rows with SELECT … FOR UPDATE SKIP LOCKED. Enqueue becomes atomic with the business transaction, which BullMQ could not have guaranteed."],
    [("Socket.IO", True, DANGER), ("Spring WebSocket + STOMP, Redis relay", True, TEXT),
     "Rooms map to destinations. Subscriptions authorised in a ChannelInterceptor using the same scope rules as REST — a Developer cannot subscribe to a ticket they could not GET."],
    [("SheetJS", True, DANGER), ("Apache POI — SXSSF write, SAX read", True, TEXT),
     "The streaming reader is required, not optional: POI's DOM reader loads a 5,000-row workbook fully into memory per concurrent import."],
    [("Prisma / TypeORM", True, DANGER), ("Spring Data JPA + Flyway", True, TEXT),
     "Migrations are hand-written SQL, which we want anyway — the immutability triggers are not expressible through a migration generator."],
    [("Shared Zod schemas", True, DANGER), ("OpenAPI → orval codegen", True, TEXT),
     "The one genuine loss. Bean Validation on the Java DTOs is the single source of truth; springdoc emits it; codegen produces Zod. CI fails on a stale client, so drift is caught at merge."],
]
table(s, ML, y, CW, ["Blueprint says", "We use", "Why, and what changes"], rows,
      [2.0, 3.3, 6.63], row_h=0.68, head_h=0.4, fsize=10.5)
pg(s)

# ---------------------------------------------------------------- 6 mysql mapping
s = slide()
y = header(s, "PostgreSQL → MySQL · type and construct mapping",
           "Eight constructs do not exist in MySQL 8 — this mapping is normative for every migration")
rows = [
    [("BIGSERIAL PRIMARY KEY", False, TEXT), ("BIGINT AUTO_INCREMENT PRIMARY KEY", False, TEXT), ""],
    [("TIMESTAMPTZ", False, TEXT), ("DATETIME(6), stored UTC", False, TEXT), ("TIMESTAMP has a 2038 limit and silent session-timezone conversion", False, TEXT2)],
    [("VARCHAR(20)[] · BIGINT[]", False, TEXT), ("JSON", False, TEXT), ("Validated in the service layer; queried with JSON_CONTAINS", False, TEXT2)],
    [("Partial index WHERE …", False, TEXT), ("Stored generated column + plain index", False, TEXT), ("pcd_open, current_ticket_id", False, TEXT2)],
    [("UPDATE … RETURNING", False, TEXT), ("LAST_INSERT_ID(expr) idiom", False, TEXT), ("Atomic, connection-local, no explicit row lock needed", False, TEXT2)],
    [("RAISE EXCEPTION in a function", False, TEXT), ("SIGNAL SQLSTATE '45000' per trigger", False, TEXT), ("MySQL needs TWO triggers where Postgres allowed one", False, TEXT2)],
    [("CHAR(64) hashes", False, TEXT), ("CHAR(64) COLLATE ascii_bin", False, TEXT), ("Exact comparison, no utf8mb4 overhead", False, TEXT2)],
    [("Boolean · NUMERIC(p,s)", False, TEXT), ("TINYINT(1) · DECIMAL(p,s)", False, TEXT), "" ],
]
table(s, ML, y, CW, ["Blueprint (PostgreSQL)", "MySQL 8 equivalent", "Note"], rows,
      [3.5, 3.9, 4.53], row_h=0.38, head_h=0.4, fsize=11)
card(s, ML, y+3.55, CW, 0.92, fill=PRIMARY_SOFT, line=None)
txt(s, ML+0.35, y+3.73, CW-0.7, 0.6,
    "Ticket ID generation:   UPDATE projects SET ticket_seq = LAST_INSERT_ID(ticket_seq + 1) WHERE id = ?;  then  SELECT LAST_INSERT_ID();\n"
    "COUNT(*) must never be used — it breaks silently, and only under concurrency.",
    size=12, color=PRIMARY, line_spacing=1.3, font=MONO)
pg(s)

# ---------------------------------------------------------------- 7 defects
s = slide()
y = header(s, "Four defects found in the blueprint itself",
           "These are in the source document, not artefacts of the MySQL move")
items = [
    ("The effort roll-up double-counts after a reopen", DANGER,
     "§4A.5's LEFT JOIN matches effort logs on (ticket_id, stage_code, iteration_no) but omits cycle_no. A ticket that "
     "re-enters the same stage at the same iteration in cycle 2 absorbs cycle 1's hours.",
     "Add  AND e.cycle_no = t.cycle_no  to the join."),
    ("The same query will not run on MySQL at all", WARNING,
     "It selects five non-aggregated columns while grouping by three. PostgreSQL permits it via functional dependency; "
     "MySQL's ONLY_FULL_GROUP_BY rejects it.",
     "Rewrite with every non-aggregated column in GROUP BY — do not disable the mode, it catches real bugs."),
    ("ticket_stage_transitions is under-protected", WARNING,
     "§4A.5 creates only a BEFORE DELETE trigger and mentions an unspecified CHECK for exited_at — leaving every other "
     "column freely updatable by a stray query or ORM dirty-check.",
     "A BEFORE UPDATE trigger permitting only the seal, using NULL-safe <=> comparison."),
    ("The hash chain forks under concurrency", DANGER,
     "§4.2 specifies prev_hash/row_hash and nightly verification but never says how two concurrent appends stay linear. "
     "They don't — both read the same latest row, and the verifier reports our own bug as tampering.",
     "Chain per-ticket, behind SELECT … FOR UPDATE on the ticket row."),
]
cy = y
for t, col, body, fix in items:
    card(s, ML, cy, CW, 1.02)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML), Inches(cy+0.1), Inches(0.055), Inches(0.82))
    bar.fill.solid(); bar.fill.fore_color.rgb = col
    bar.line.fill.background(); bar.shadow.inherit = False
    txt(s, ML+0.3, cy+0.12, 4.5, 0.5, t, size=13, bold=True, color=TEXT, line_spacing=1.15)
    txt(s, ML+0.3, cy+0.62, 4.5, 0.32, "→  " + fix, size=10, bold=True, color=col, line_spacing=1.15)
    txt(s, ML+5.1, cy+0.16, CW-5.4, 0.75, body, size=11, color=TEXT2, line_spacing=1.26)
    cy += 1.13
pg(s)

# ---------------------------------------------------------------- 8 immutability
s = slide()
y = header(s, "Immutability in MySQL", "Triggers, their three limits, and what closes each")
code_block(s, ML, y, 6.6, 2.5,
           "CREATE TRIGGER trg_hist_no_update BEFORE UPDATE ON ticket_history\n"
           "FOR EACH ROW\n"
           "BEGIN\n"
           "  SIGNAL SQLSTATE '45000'\n"
           "    SET MESSAGE_TEXT = 'Immutable table: ticket_history\n"
           "                        rows cannot be updated';\n"
           "END\n"
           "\n"
           "-- and a second, identical trigger for BEFORE DELETE\n"
           "-- MySQL cannot express  BEFORE UPDATE OR DELETE", size=11)
x2 = ML + 6.9
lims = [
    ("Triggers do not fire on TRUNCATE", "The app DB user is granted no DROP privilege, which TRUNCATE requires"),
    ("Triggers do not survive DROP TABLE", "No DDL privileges for the app user; Flyway runs as a separate migration user"),
    ("A privileged user can drop the triggers", "The hash chain is the backstop — tampering that bypasses triggers still breaks the chain"),
]
cy = y
for t, f in lims:
    card(s, x2, cy, CW-6.9, 0.78)
    txt(s, x2+0.26, cy+0.11, CW-7.4, 0.28, t, size=11.5, bold=True, color=DANGER)
    txt(s, x2+0.26, cy+0.4, CW-7.4, 0.32, f, size=10.5, color=TEXT2, line_spacing=1.2)
    cy += 0.87
card(s, ML, y+2.8, CW, 1.65, fill=RGBColor(0x1E,0x1B,0x33), line=None)
txt(s, ML+0.35, y+3.0, CW-0.7, 0.3, "THE ONE PERMITTED MUTATION", size=10.5, bold=True, color=RGBColor(0xA5,0xB4,0xFC))
txt(s, ML+0.35, y+3.38, CW-0.7, 0.95,
    "Sealing a stage transition: exited_at goes NULL → timestamp, and duration_mins is computed at the same moment.\n"
    "The trigger permits it only when the row is still open AND only those columns changed — compared with <=>, "
    "because plain = returns NULL for two NULLs and silently passes the check.",
    size=12.5, color=WHITE, line_spacing=1.3)
pg(s)

# ---------------------------------------------------------------- 9 deviations
s = slide()
y = header(s, "Every deviation from the blueprint, in one place", "Nothing is silent")
rows = [
    [("D-1", True, PRIMARY), "NestJS / Node backend", "Spring Boot / Java", ("Directed by the team", False, TEXT2)],
    [("D-2", True, PRIMARY), "PostgreSQL 16", "MySQL 8.4", ("Directed by the team", False, TEXT2)],
    [("D-3", True, PRIMARY), "BullMQ, Socket.IO, SheetJS, Prisma", "Outbox, STOMP, POI, JPA+Flyway", ("No Java equivalents", False, TEXT2)],
    [("D-4", True, PRIMARY), "Shared Zod schemas", "OpenAPI + generated client, CI check", ("Cannot share validation across languages", False, TEXT2)],
    [("D-5", True, DANGER), "DELETE trigger + unspecified rule", "Explicit BEFORE UPDATE seal-only trigger", ("Every other column was left updatable", False, TEXT2)],
    [("D-6", True, DANGER), "Hash chain, concurrency unspecified", "Per-ticket chain with row lock", ("Concurrent appends fork the chain", False, TEXT2)],
    [("D-7", True, DANGER), "Roll-up query as written", "Rewritten + cycle_no added to join", ("Original double-counts after a reopen", False, TEXT2)],
    [("D-8", True, WARNING), "Ticket ID year behaviour unstated", "No reset at rollover; code immutable", ("Resetting creates a midnight collision window", False, TEXT2)],
    [("D-9", True, WARNING), "Phases 0–6 by calendar", "Milestones M0–M7 by dependency", ("Each milestone independently demoable", False, TEXT2)],
]
table(s, ML, y, CW, ["#", "Blueprint", "What we do instead", "Why"], rows,
      [0.7, 3.4, 3.9, 3.93], row_h=0.38, head_h=0.4, fsize=10.5)
pg(s)

# ---------------------------------------------------------------- 10 decisions
s = slide()
y = header(s, "Four decisions must be locked before M2",
           "Blueprint §16 lists 13 governance questions — nine can wait, these four change the schema")
rows = [
    [("G-1", True, PRIMARY), "Is effort mandatory at handoff — blocking or warn-only?",
     ("Blocking, per-project flag defaulting on", True, TEXT), ("Handoff DTO, S-29, skip report", False, TEXT2)],
    [("G-2", True, PRIMARY), "Does a rework loop reset the Planned Close Date?",
     ("No — the original date stands", True, TEXT), ("SLA scanner, breach reporting", False, TEXT2)],
    [("G-3", True, PRIMARY), "Can a Developer close, or only mark Resolved?",
     ("Resolved only", True, TEXT), ("workflow_transitions seed, S-23", False, TEXT2)],
    [("G-4", True, PRIMARY), "Does an auto-escalated level revert after closure?",
     ("No — original_level preserved, both reported", True, TEXT), ("Already in schema; needs the report", False, TEXT2)],
]
table(s, ML, y, CW, ["#", "Question", "Proposed default", "What it touches"], rows,
      [0.7, 4.4, 3.9, 2.93], row_h=0.6, head_h=0.4, fsize=11)
card(s, ML, y+3.0, CW, 1.45, fill=RGBColor(0xFE,0xF3,0xC7), line=None)
txt(s, ML+0.35, y+3.2, CW-0.7, 0.3, "OPEN QUESTION THE BLUEPRINT DOES NOT ANSWER", size=10.5, bold=True, color=RGBColor(0x92,0x40,0x0E))
txt(s, ML+0.35, y+3.58, CW-0.7, 0.75,
    "Is there an existing employee directory or SSO — Azure AD, Google Workspace — that the Resource Master should sync from, "
    "or is EduTrack the system of record for users?\n"
    "If a directory exists, users needs an external ID from M1, not after M3.",
    size=12.5, color=RGBColor(0x92,0x40,0x0E), line_spacing=1.3)
pg(s)

# ---------------------------------------------------------------- 11 milestones
s = slide()
y = header(s, "Milestones — sequenced by dependency, not calendar",
           "Each ends in something demoable, with explicit exit criteria")
rows = [
    [("M0", True, A_C), ("Foundation", True, TEXT), "Monorepo, Docker Compose, the complete schema in one pass, triggers, seeds, design tokens, CI",
     ("docker compose up gives a migrated, seeded DB", False, TEXT2)],
    [("M1", True, A_C), ("Auth + scope guard", True, TEXT), "Argon2id login, JWT + refresh rotation with family revocation, the three guards, permission matrix",
     ("All six roles land correctly; Developer gets 404", False, TEXT2)],
    [("M2", True, A_C), ("Immutability core", True, TEXT), "Insert-only services, canonical JSON, per-ticket hash chain, nightly verifier, negative tests",
     ("Verifier passes; no fork under concurrent load", False, TEXT2)],
    [("M3", True, B_C), ("Masters", True, TEXT), "13 master screens, client master, the import wizard built once and registered twice, working calendar",
     ("Admin can stand up a tenant without touching the DB", False, TEXT2)],
    [("M4", True, C_C), ("Tickets + ribbon", True, TEXT), "Ticket CRUD, cycles, reopen, comments, attachments, transition service, ribbon, Journey grid",
     ("Walkthrough A reconciles to 38.0 h, 3 iterations", False, TEXT2)],
    [("M5", True, D_C), ("Engines", True, TEXT), "SLA + stage-SLA scanners on the working calendar, escalation, mail engine, notification centre",
     ("Scanner escalates, mails, turns the banner red live", False, TEXT2)],
    [("M6", True, A_C), ("Dashboard + reports", True, TEXT), "Pre-aggregated summary tables, 20 widgets with drill-down, 18 reports with export",
     ("First paint under 1.5 s at 50,000 tickets", False, TEXT2)],
    [("M7", True, D_C), ("Chat + hardening", True, TEXT), "STOMP chat, Ask Status, then CSP, scanning, load test, penetration test, UAT",
     ("Security review passes; UAT signed off", False, TEXT2)],
]
table(s, ML, y, CW, ["", "Milestone", "Scope", "Exit criterion"], rows,
      [0.62, 2.1, 5.3, 3.91], row_h=0.44, head_h=0.4, fsize=10)
pg(s)

# ---------------------------------------------------------------- 12 testing
s = slide()
y = header(s, "Testing strategy",
           "Weighted toward the three risks the blueprint identifies as most expensive")
rows = [
    [("Unit", True, TEXT), "JUnit 5 + AssertJ", "SLA maths on the working calendar, effort roll-ups, iteration/cycle counters, hash canonicalisation"],
    [("Integration", True, TEXT), "Testcontainers", "Every migration against real MySQL; triggers proven by negative tests; reopen and handoff transactions"],
    [("Permission matrix", True, DANGER), "Parameterised MockMvc", "Every role × every route. A new route without a matrix entry fails the build"],
    [("Architecture", True, TEXT), "ArchUnit", "No controller reaches a repository directly; no update/delete method exists on append-only services"],
    [("Concurrency", True, DANGER), "Custom harness", "Parallel ticket creation (no duplicate codes); parallel appends (no chain fork); parallel handoffs"],
    [("Contract", True, TEXT), "springdoc + codegen diff", "CI fails if the committed TypeScript client is stale against the spec"],
    [("End-to-end", True, TEXT), "Playwright", "Walkthroughs A, B and C from blueprint §14, against a seeded database"],
    [("Performance", True, TEXT), "k6", "Dashboard and ticket-list p95 on a 50,000-ticket dataset"],
]
table(s, ML, y, CW, ["Layer", "Tooling", "What it must cover"], rows,
      [2.2, 2.5, 7.23], row_h=0.42, head_h=0.4, fsize=11)
pg(s)

# ---------------------------------------------------------------- 13 risks
s = slide()
y = header(s, "Risks carried, and four specific to this stack", "")
rows = [
    [("MySQL trigger protection bypassed by a privileged user", True, TEXT),
     "The hash chain is the backstop; nightly verification; the app user holds no DDL privilege"],
    [("Generated-column indexes underperform Postgres partial indexes as data grows", True, TEXT),
     "Benchmark the SLA scan at M5 against 500,000 rows; partition tickets by year if needed"],
    [("OpenAPI client drift between Java and TypeScript", True, TEXT),
     "CI staleness check fails the build; Bean Validation is the single source of truth"],
    [("Outbox worker falls behind under mail burst", True, TEXT),
     "SELECT … FOR UPDATE SKIP LOCKED allows horizontal worker scaling; email_log depth is a Prometheus alert"],
]
table(s, ML, y, CW, ["Stack-specific risk", "Mitigation"], rows,
      [5.6, 6.33], row_h=0.6, head_h=0.4, fsize=11.5)
card(s, ML, y+2.85, CW, 1.6, fill=PRIMARY_SOFT, line=None)
txt(s, ML+0.35, y+3.05, CW-0.7, 0.3, "PLUS THE BLUEPRINT'S OWN REGISTER, UNCHANGED", size=10.5, bold=True, color=PRIMARY)
txt(s, ML+0.35, y+3.45, CW-0.7, 0.85,
    "History integrity challenged in a dispute · notification fatigue · dashboard slowdown · under-logged effort · "
    "role-scope leakage · ribbon unreadable at 8 stages · teams gaming the stage clock · missed alerts · "
    "client import corruption · malicious uploads · internal notes leaking to a client.",
    size=12.5, color=PRIMARY, line_spacing=1.32)
pg(s)

# ---------------------------------------------------------------- 14 next
s = slide(SURFACE); grad(s)
txt(s, ML+0.3, 1.0, CW, 0.55, "What happens next", size=32, bold=True, color=WHITE)
ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML+0.3), Inches(1.78), Inches(1.3), Inches(0.04))
ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0xA5,0xB4,0xFC)
ln.line.fill.background(); ln.shadow.inherit = False
card(s, ML+0.3, 2.25, CW-0.6, 1.5, fill=RGBColor(0xFF,0xFF,0xFF), line=None)
txt(s, ML+0.65, 2.5, CW-1.3, 0.35, "IMMEDIATE NEXT STEP  ·  M0", size=11, bold=True, color=PRIMARY)
txt(s, ML+0.65, 2.92, CW-1.3, 0.7,
    "Repo scaffold, Docker Compose, the complete Flyway migration set translated per §3, seed data, and the design "
    "tokens — a running, migrated stack to inspect before any feature work begins.",
    size=14, color=TEXT, line_spacing=1.3)
txt(s, ML+0.3, 4.15, CW, 0.35, "BLOCKING ON NOTHING", size=11, bold=True, color=RGBColor(0xA5,0xB4,0xFC))
txt(s, ML+0.3, 4.55, CW-0.6, 0.9,
    "The four governance decisions are needed before M2, not before M0.\n"
    "The SSO / directory question is needed before M1.",
    size=15, color=WHITE, line_spacing=1.35)
txt(s, ML+0.3, 6.6, CW, 0.3,
    "PLAN.md  ·  TEAM-PLAN.md  ·  GETTING-STARTED.md  ·  docs/streams/STREAM-A…D.md",
    size=11.5, color=RGBColor(0x81,0x8C,0xF8))

save(str(pathlib.Path(__file__).parents[2] / "docs/decks") + "/EduTrack-Implementation-Plan.pptx")
