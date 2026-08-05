from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------- tokens
# from blueprint 12.1
PRIMARY      = RGBColor(0x4F, 0x46, 0xE5)
PRIMARY_SOFT = RGBColor(0xEE, 0xF2, 0xFF)
BG_APP       = RGBColor(0xF7, 0xF8, 0xFC)
SURFACE      = RGBColor(0xFF, 0xFF, 0xFF)
SUBTLE       = RGBColor(0xF1, 0xF3, 0xF9)
BORDER       = RGBColor(0xE5, 0xE8, 0xF0)
TEXT         = RGBColor(0x11, 0x18, 0x27)
TEXT2        = RGBColor(0x6B, 0x72, 0x80)
SUCCESS      = RGBColor(0x10, 0xB9, 0x81)
WARNING      = RGBColor(0xF5, 0x9E, 0x0B)
DANGER       = RGBColor(0xEF, 0x44, 0x44)
INFO         = RGBColor(0x3B, 0x82, 0xF6)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)

# stream accents
A_C = RGBColor(0x4F, 0x46, 0xE5)
B_C = RGBColor(0x06, 0xB6, 0xD4)
C_C = RGBColor(0x8B, 0x5C, 0xF6)
D_C = RGBColor(0x14, 0xB8, 0xA6)

FONT = "Calibri"
MONO = "Consolas"

SW, SH = 13.333, 7.5
ML = 0.7                      # left margin
CW = SW - 2 * ML              # content width

prs = None
BLANK = None


def new_deck():
    global prs, BLANK
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    BLANK = prs.slide_layouts[6]
    return prs


def save(path):
    prs.save(path)
    print("saved:", path, "| slides:", len(prs.slides._sldIdLst))


# ---------------------------------------------------------------- helpers
def slide(bg=BG_APP):
    s = prs.slides.add_slide(BLANK)
    bgshape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0,
                                 Inches(SW), Inches(SH))
    bgshape.fill.solid()
    bgshape.fill.fore_color.rgb = bg
    bgshape.line.fill.background()
    bgshape.shadow.inherit = False
    return s


def txt(s, x, y, w, h, text, size=14, bold=False, color=TEXT,
        align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP, space_after=0,
        line_spacing=1.15):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = font
    return tb


def header(s, title, sub=None, accent=PRIMARY):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(ML), Inches(0.52),
                             Inches(0.055), Inches(0.46))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    bar.shadow.inherit = False
    txt(s, ML + 0.2, 0.5, CW - 0.2, 0.5, title, size=27, bold=True, color=TEXT)
    if sub:
        txt(s, ML + 0.2, 1.05, CW - 0.2, 0.36, sub, size=13.5, color=TEXT2)
        return 1.62
    return 1.28


def card(s, x, y, w, h, fill=SURFACE, line=BORDER, radius=True):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    shp.shadow.inherit = False
    if radius:
        try:
            shp.adjustments[0] = 0.06
        except Exception:
            pass
    return shp


def table(s, x, y, w, cols, rows, col_w, header_fill=PRIMARY,
          header_color=WHITE, fsize=11.5, hsize=11.5, row_h=0.34,
          head_h=0.4, zebra=True):
    """rows: list of list of (text, bold, color) or plain str"""
    n_rows = len(rows) + 1
    shape = s.shapes.add_table(n_rows, len(cols), Inches(x), Inches(y),
                               Inches(w), Inches(head_h + row_h * len(rows)))
    tbl = shape.table
    tbl.first_row = False
    tbl.horz_banding = False
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = Inches(cw)
    tbl.rows[0].height = Inches(head_h)
    for r in range(1, n_rows):
        tbl.rows[r].height = Inches(row_h)

    def fill_cell(cell, value, bold, color, fill_rgb, size):
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill_rgb
        cell.margin_left = Inches(0.09)
        cell.margin_right = Inches(0.06)
        cell.margin_top = Inches(0.03)
        cell.margin_bottom = Inches(0.03)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.line_spacing = 1.0
        r = p.add_run()
        r.text = value
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = FONT

    for c, head in enumerate(cols):
        fill_cell(tbl.cell(0, c), head, True, header_color, header_fill, hsize)
    for ri, row in enumerate(rows):
        bg = SURFACE if (not zebra or ri % 2 == 0) else SUBTLE
        for ci, val in enumerate(row):
            if isinstance(val, tuple):
                t, b, col = val
            else:
                t, b, col = val, False, TEXT
            fill_cell(tbl.cell(ri + 1, ci), t, b, col, bg, fsize)
    return shape


def bullets(s, x, y, w, items, size=14, gap=0.42, dot=PRIMARY, bold_lead=True):
    """items: list of (lead, rest) or plain string"""
    cy = y
    for it in items:
        d = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(cy + 0.075),
                               Inches(0.085), Inches(0.085))
        d.fill.solid()
        d.fill.fore_color.rgb = dot
        d.line.fill.background()
        d.shadow.inherit = False
        tb = s.shapes.add_textbox(Inches(x + 0.22), Inches(cy),
                                  Inches(w - 0.22), Inches(gap))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.line_spacing = 1.2
        if isinstance(it, tuple):
            lead, rest = it
            r1 = p.add_run(); r1.text = lead
            r1.font.size = Pt(size); r1.font.bold = bold_lead
            r1.font.color.rgb = TEXT; r1.font.name = FONT
            r2 = p.add_run(); r2.text = rest
            r2.font.size = Pt(size); r2.font.color.rgb = TEXT2
            r2.font.name = FONT
        else:
            r = p.add_run(); r.text = it
            r.font.size = Pt(size); r.font.color.rgb = TEXT
            r.font.name = FONT
        cy += gap
    return cy


def footer(s, n, label="EduTrack"):
    txt(s, ML, SH - 0.46, 4.0, 0.26, label, size=9.5, color=TEXT2)
    txt(s, SW - ML - 1.2, SH - 0.46, 1.2, 0.26, str(n), size=9.5,
        color=TEXT2, align=PP_ALIGN.RIGHT)


def code_block(s, x, y, w, h, lines, size=11.5):
    card(s, x, y, w, h, fill=RGBColor(0x1E, 0x1B, 0x33), line=None)
    txt(s, x + 0.22, y + 0.16, w - 0.44, h - 0.32, lines, size=size,
        color=RGBColor(0xE5, 0xE8, 0xF0), font=MONO, line_spacing=1.28)


